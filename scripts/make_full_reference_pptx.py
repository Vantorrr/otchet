#!/usr/bin/env python3
"""Generate FULL reference-style PPTX with gradients, shadows, and effects."""
from __future__ import annotations

import os
import sys
from datetime import datetime

from dotenv import load_dotenv

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from bot.config import Settings
from bot.services.di import Container
from bot.services.data_aggregator import DataAggregatorService
from bot.services.presentation import PresentationService

# PPTX advanced imports
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.oxml.shapes.shared import qn
from pptx.oxml import parse_xml


def add_gradient_fill(shape, color1: str, color2: str):
    """Add gradient fill to shape."""
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = hex_to_rgb(color1)
    fill.gradient_stops[1].color.rgb = hex_to_rgb(color2)


def add_shadow(shape):
    """Add shadow effect to shape."""
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.distance = Pt(3)
        shadow.blur_radius = Pt(6)
        shadow.color.rgb = RGBColor(0, 0, 0)
        shadow.alpha = 0.3
    except Exception:
        pass


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def main():
    load_dotenv()
    os.environ.setdefault("GOOGLE_APPLICATION_CREDENTIALS", "service_account.json")
    os.environ.setdefault("BOT_TOKEN", "dummy-token-for-offline")

    settings = Settings.load()
    Container.init(settings)

    start = datetime.strptime("2025-09-01", "%Y-%m-%d").date()
    end = datetime.strptime("2025-09-07", "%Y-%m-%d").date()
    print(f"Building FULL REFERENCE PPTX with gradients/shadows for {start}..{end}")

    aggregator = DataAggregatorService(Container.get().sheets)
    period_data, prev_data, period_name, start_date, end_date, prev_start, prev_end = (
        __import__("asyncio").get_event_loop().run_until_complete(
            aggregator.aggregate_custom_with_previous(start, end)
        )
    )
    if not period_data:
        print("No data for the requested period.")
        return 1

    prs_service = PresentationService(settings)
    totals = prs_service._calculate_totals(period_data)
    
    # Create PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # === TITLE SLIDE WITH GRADIENT ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Gradient background (emerald to dark emerald)
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    add_gradient_fill(bg_shape, "#1B5E20", "#2E7D32")
    bg_shape.line.fill.background()
    
    # Geometric decorative elements
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(2), Inches(2))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = hex_to_rgb("#4CAF50")
    circle1.line.fill.background()
    circle1.fill.transparency = 0.3
    
    # Large decorative circle
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4.5), Inches(2.5), Inches(2.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    circle2.line.fill.background()
    circle2.fill.transparency = 0.8
    
    # Triangle accent
    triangle = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(1), Inches(5), Inches(1.5), Inches(1.5))
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = hex_to_rgb("#66BB6A")
    triangle.line.fill.background()
    triangle.fill.transparency = 0.4
    
    # Main title with shadow
    title_box = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(7.33), Inches(2.5))
    title_frame = title_box.text_frame
    title_frame.text = f"{getattr(settings, 'office_name', 'БАНКОВСКИЕ ГАРАНТИИ').upper()}\n\nОТЧЕТ ПО ПРОДАЖАМ"
    for p in title_frame.paragraphs:
        p.font.size = Pt(48)
        p.font.name = "Roboto"
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    add_shadow(title_box)
    
    # Period subtitle
    period_box = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7.33), Inches(1.5))
    period_frame = period_box.text_frame
    period_frame.text = f"{period_name}\n{start.strftime('%d.%m.%Y')} — {end.strftime('%d.%m.%Y')}"
    for i, p in enumerate(period_frame.paragraphs):
        p.font.size = Pt(28 if i == 0 else 20)
        p.font.name = "Roboto"
        p.font.color.rgb = hex_to_rgb("#E8F5E8")
        p.alignment = PP_ALIGN.CENTER
    
    # === DASHBOARD WITH CARDS AND SHADOWS ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background (light to cream)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    add_gradient_fill(bg2, "#FAFAFA", "#F5F5DC")
    bg2.line.fill.background()
    
    # Header with gradient
    header_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    add_gradient_fill(header_bar, "#2E7D32", "#1B5E20")
    header_bar.line.fill.background()
    
    # Header text with glow effect
    header_text = slide2.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.33), Inches(0.6))
    header_text.text_frame.text = "КЛЮЧЕВЫЕ ПОКАЗАТЕЛИ КОМАНДЫ"
    header_text.text_frame.paragraphs[0].font.size = Pt(26)
    header_text.text_frame.paragraphs[0].font.name = "Roboto"
    header_text.text_frame.paragraphs[0].font.bold = True
    header_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Metrics cards with shadows and gradients
    metrics = [
        ("ПОВТОРНЫЕ ЗВОНКИ", f"{int(totals.get('calls_fact', 0)):,}".replace(",", " "), f"{totals.get('calls_percentage', 0):.1f}%".replace(".", ","), "#4CAF50"),
        ("ЗАЯВКИ (ШТ)", f"{int(totals.get('leads_units_fact', 0)):,}".replace(",", " "), f"{totals.get('leads_units_percentage', 0):.1f}%".replace(".", ","), "#66BB6A"),
        ("ЗАЯВКИ (МЛН)", f"{totals.get('leads_volume_fact', 0):.1f}".replace(".", ","), f"{totals.get('leads_volume_percentage', 0):.1f}%".replace(".", ","), "#81C784"),
        ("ОДОБРЕНО (МЛН)", f"{totals.get('approved_volume', 0):.1f}".replace(".", ","), "—", "#A5D6A7"),
        ("ВЫДАНО (МЛН)", f"{totals.get('issued_volume', 0):.1f}".replace(".", ","), "—", "#C8E6C9"),
        ("НОВЫЕ ЗВОНКИ", f"{int(totals.get('new_calls', 0)):,}".replace(",", " "), "—", "#E1F5FE"),
    ]
    
    card_w = Inches(4)
    card_h = Inches(1.8)
    for i, (label, value, pct, accent_color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.2)
        y = Inches(1.7 + row * 2.1)
        
        # Card with gradient and shadow
        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        add_gradient_fill(card, "#FFFFFF", "#F8F9FA")
        card.line.color.rgb = hex_to_rgb("#E0E0E0")
        card.line.width = Pt(0.5)
        add_shadow(card)
        
        # Accent top stripe with gradient
        accent = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.15))
        add_gradient_fill(accent, accent_color, "#2E7D32")
        accent.line.fill.background()
        
        # Icon background (circular)
        icon_bg = slide2.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.3), Inches(0.6), Inches(0.6))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = hex_to_rgb(accent_color)
        icon_bg.line.fill.background()
        icon_bg.fill.transparency = 0.2
        
        # Label
        label_box = slide2.shapes.add_textbox(x + Inches(1), y + Inches(0.25), card_w - Inches(1.2), Inches(0.5))
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(12)
        label_box.text_frame.paragraphs[0].font.name = "Roboto"
        label_box.text_frame.paragraphs[0].font.bold = True
        label_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#424242")
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Large value
        value_box = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), Inches(0.7))
        value_box.text_frame.text = value
        value_box.text_frame.paragraphs[0].font.size = Pt(32)
        value_box.text_frame.paragraphs[0].font.name = "Roboto"
        value_box.text_frame.paragraphs[0].font.bold = True
        value_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#1B5E20")
        value_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Percentage badge
        if pct != "—":
            pct_shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + card_w - Inches(1.2), y + Inches(0.2), Inches(1), Inches(0.4))
            pct_shape.fill.solid()
            pct_shape.fill.fore_color.rgb = hex_to_rgb(accent_color)
            pct_shape.line.fill.background()
            
            pct_text = slide2.shapes.add_textbox(x + card_w - Inches(1.15), y + Inches(0.3), Inches(0.9), Inches(0.2))
            pct_text.text_frame.text = pct
            pct_text.text_frame.paragraphs[0].font.size = Pt(11)
            pct_text.text_frame.paragraphs[0].font.name = "Roboto"
            pct_text.text_frame.paragraphs[0].font.bold = True
            pct_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            pct_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # === RANKING WITH LARGE VISUAL CARDS ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient background
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    add_gradient_fill(bg3, "#E8F5E8", "#F1F8E9")
    bg3.line.fill.background()
    
    # Header with emerald gradient
    header3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    add_gradient_fill(header3, "#1B5E20", "#2E7D32")
    header3.line.fill.background()
    
    # Header text
    header3_text = slide3.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11.33), Inches(0.6))
    header3_text.text_frame.text = "РЕЙТИНГ ЭФФЕКТИВНОСТИ"
    header3_text.text_frame.paragraphs[0].font.size = Pt(28)
    header3_text.text_frame.paragraphs[0].font.name = "Roboto"
    header3_text.text_frame.paragraphs[0].font.bold = True
    header3_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header3_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Build ranking
    scored = []
    for m in period_data.values():
        calls_pct = (m.calls_fact/m.calls_plan*100) if m.calls_plan else 0
        vol_pct = (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0
        scored.append((0.5*calls_pct+0.5*vol_pct, m.name))
    scored.sort(reverse=True)
    best = [n for _, n in scored[:2]]
    worst = [n for _, n in list(reversed(scored[-2:]))]
    
    # BEST PERFORMERS - Large card with gradient and shadow
    best_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(5.5), Inches(4.5))
    add_gradient_fill(best_card, "#2E7D32", "#1B5E20")
    best_card.line.fill.background()
    add_shadow(best_card)
    
    # Trophy icon with glow
    trophy_bg = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.3), Inches(1.2), Inches(1.2))
    trophy_bg.fill.solid()
    trophy_bg.fill.fore_color.rgb = hex_to_rgb("#FFD700")
    trophy_bg.line.fill.background()
    add_shadow(trophy_bg)
    
    trophy_text = slide3.shapes.add_textbox(Inches(1.8), Inches(2.6), Inches(0.6), Inches(0.6))
    trophy_text.text_frame.text = "👑"
    trophy_text.text_frame.paragraphs[0].font.size = Pt(32)
    trophy_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Best title
    best_title = slide3.shapes.add_textbox(Inches(3), Inches(2.4), Inches(3.2), Inches(0.8))
    best_title.text_frame.text = "ЛИДЕРЫ ПЕРИОДА"
    best_title.text_frame.paragraphs[0].font.size = Pt(22)
    best_title.text_frame.paragraphs[0].font.name = "Roboto"
    best_title.text_frame.paragraphs[0].font.bold = True
    best_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Best performers list
    for i, name in enumerate(best[:2]):
        name_box = slide3.shapes.add_textbox(Inches(1.3), Inches(3.5 + i*1), Inches(4.9), Inches(0.7))
        name_box.text_frame.text = f"🏆 {i+1}. {name}"
        name_box.text_frame.paragraphs[0].font.size = Pt(18)
        name_box.text_frame.paragraphs[0].font.name = "Roboto"
        name_box.text_frame.paragraphs[0].font.bold = True
        name_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        desc_box = slide3.shapes.add_textbox(Inches(1.3), Inches(3.8 + i*1), Inches(4.9), Inches(0.4))
        desc_box.text_frame.text = "превышение плана по ключевым показателям"
        desc_box.text_frame.paragraphs[0].font.size = Pt(12)
        desc_box.text_frame.paragraphs[0].font.name = "Roboto"
        desc_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#C8E6C9")
    
    # UNDERPERFORMERS - Red gradient card
    worst_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5))
    add_gradient_fill(worst_card, "#D32F2F", "#B71C1C")
    worst_card.line.fill.background()
    add_shadow(worst_card)
    
    # Warning icon
    warn_bg = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(2.3), Inches(1.2), Inches(1.2))
    warn_bg.fill.solid()
    warn_bg.fill.fore_color.rgb = hex_to_rgb("#FF9800")
    warn_bg.line.fill.background()
    add_shadow(warn_bg)
    
    warn_text = slide3.shapes.add_textbox(Inches(7.8), Inches(2.6), Inches(0.6), Inches(0.6))
    warn_text.text_frame.text = "⚠️"
    warn_text.text_frame.paragraphs[0].font.size = Pt(32)
    warn_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Worst title
    worst_title = slide3.shapes.add_textbox(Inches(9), Inches(2.4), Inches(3.2), Inches(0.8))
    worst_title.text_frame.text = "ТРЕБУЮТ ВНИМАНИЯ"
    worst_title.text_frame.paragraphs[0].font.size = Pt(22)
    worst_title.text_frame.paragraphs[0].font.name = "Roboto"
    worst_title.text_frame.paragraphs[0].font.bold = True
    worst_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Underperformers list
    for i, name in enumerate(worst[:2]):
        name_box = slide3.shapes.add_textbox(Inches(7.3), Inches(3.5 + i*1), Inches(4.9), Inches(0.7))
        name_box.text_frame.text = f"📉 {i+1}. {name}"
        name_box.text_frame.paragraphs[0].font.size = Pt(18)
        name_box.text_frame.paragraphs[0].font.name = "Roboto"
        name_box.text_frame.paragraphs[0].font.bold = True
        name_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        desc_box = slide3.shapes.add_textbox(Inches(7.3), Inches(3.8 + i*1), Inches(4.9), Inches(0.4))
        desc_box.text_frame.text = "отставание от плановых показателей"
        desc_box.text_frame.paragraphs[0].font.size = Pt(12)
        desc_box.text_frame.paragraphs[0].font.name = "Roboto"
        desc_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#FFCDD2")
    
    # Save
    out = f"full_reference_{start.strftime('%Y%m%d')}_{end.strftime('%Y%m%d')}.pptx"
    prs.save(out)
    print(f"FULL REFERENCE PPTX saved: {out}")
    print("This shows the target design - will port to Slides when quota is resolved.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

