#!/usr/bin/env python3
"""МОЩНАЯ презентация с ВСЕМИ диаграммами из эскизов + премиум‑эффекты."""
from __future__ import annotations

import os
import sys
from datetime import datetime

from dotenv import load_dotenv

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from bot.config import Settings
from bot.services.di import Container
from bot.services.data_aggregator import DataAggregatorService
from bot.services.presentation import PresentationService
from bot.services.yandex_gpt import YandexGPTService

# PPTX advanced + plotly
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import kaleido


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def add_gradient_fill(shape, color1: str, color2: str):
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = hex_to_rgb(color1)
    fill.gradient_stops[1].color.rgb = hex_to_rgb(color2)


def add_shadow(shape):
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.distance = Pt(4)
        shadow.blur_radius = Pt(8)
        shadow.color.rgb = RGBColor(0, 0, 0)
        shadow.alpha = 0.25
    except Exception:
        pass


def add_logo_to_slide(slide, prs, logo_path):
    """Добавить логотип справа сверху на слайд."""
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(
                logo_path,
                prs.slide_width - Inches(2.2), Inches(0.2),
                height=Inches(0.8)
            )
        except Exception:
            pass


def create_calls_chart(daily_data, output_path="calls_chart.png"):
    """График звонков план/факт по дням."""
    dates = [d['date'] for d in daily_data]
    plan = [d.get('calls_plan', 0) for d in daily_data]
    fact = [d.get('calls_fact', 0) for d in daily_data]
    
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=dates, y=plan, mode='lines+markers', name='План', 
                            line=dict(color='#4CAF50', width=3), marker=dict(size=8)))
    fig.add_trace(go.Scatter(x=dates, y=fact, mode='lines+markers', name='Факт',
                            line=dict(color='#2196F3', width=3), marker=dict(size=8)))
    
    fig.update_layout(
        title="Звонки: план vs факт", title_font_size=16, title_font_family="Roboto",
        xaxis_title="Дни", yaxis_title="Количество звонков",
        font=dict(family="Roboto", size=12), 
        plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',
        width=600, height=350
    )
    fig.write_image(output_path)
    return output_path


def create_spider_chart(manager_data, avg_data, manager_name, output_path="spider_chart.png"):
    """Spider‑диаграмма менеджер vs отдел."""
    categories = ['Повторные звонки', 'Новые звонки', 'Заявки шт', 'Заявки млн', 'Одобрено', 'Выдано']
    
    manager_values = [
        manager_data.calls_fact, manager_data.new_calls, manager_data.leads_units_fact,
        manager_data.leads_volume_fact, manager_data.approved_volume, manager_data.issued_volume
    ]
    
    avg_values = [
        avg_data.get('calls_fact', 0), avg_data.get('new_calls', 0), avg_data.get('leads_units_fact', 0),
        avg_data.get('leads_volume_fact', 0), avg_data.get('approved_volume', 0), avg_data.get('issued_volume', 0)
    ]
    
    fig = go.Figure()
    
    fig.add_trace(go.Scatterpolar(
        r=avg_values + [avg_values[0]],  # Close the polygon
        theta=categories + [categories[0]],
        fill='toself',
        name='Среднее по отделу',
        fillcolor='rgba(255, 140, 101, 0.3)',  # #FF8A65 with transparency
        line=dict(color='#FF8A65', width=2)
    ))
    
    fig.add_trace(go.Scatterpolar(
        r=manager_values + [manager_values[0]],
        theta=categories + [categories[0]],
        fill='toself',
        name=manager_name,
        fillcolor='rgba(156, 39, 176, 0.4)',  # Purple with transparency
        line=dict(color='#9C27B0', width=3)
    ))
    
    fig.update_layout(
        polar=dict(
            radialaxis=dict(visible=True, range=[0, max(max(manager_values), max(avg_values)) * 1.1])
        ),
        title=f"Сравнение — {manager_name}",
        title_font_size=16, title_font_family="Roboto",
        font=dict(family="Roboto", size=11),
        width=500, height=400,
        paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)'
    )
    fig.write_image(output_path)
    return output_path


def create_bar_chart(managers_data, output_path="managers_bar.png"):
    """Столбчатая диаграмма по менеджерам."""
    names = [m.name for m in managers_data]
    calls = [m.calls_fact for m in managers_data]
    leads = [m.leads_units_fact for m in managers_data]
    
    fig = go.Figure()
    fig.add_trace(go.Bar(name='Звонки', x=names, y=calls, marker_color='#4CAF50'))
    fig.add_trace(go.Bar(name='Заявки', x=names, y=leads, marker_color='#2196F3'))
    
    fig.update_layout(
        title="Результаты по менеджерам", title_font_size=16, title_font_family="Roboto",
        xaxis_title="Менеджеры", yaxis_title="Количество",
        font=dict(family="Roboto", size=12),
        plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)',
        width=700, height=400
    )
    fig.write_image(output_path)
    return output_path


def main():
    load_dotenv()
    os.environ.setdefault("GOOGLE_APPLICATION_CREDENTIALS", "service_account.json")
    os.environ.setdefault("BOT_TOKEN", "dummy-token-for-offline")

    settings = Settings.load()
    Container.init(settings)

    start = datetime.strptime("2025-09-01", "%Y-%m-%d").date()
    end = datetime.strptime("2025-09-07", "%Y-%m-%d").date()
    print(f"Building POWERHOUSE PPTX with ALL charts for {start}..{end}")

    aggregator = DataAggregatorService(Container.get().sheets)
    period_data, prev_data, period_name, start_date, end_date, prev_start, prev_end = (
        __import__("asyncio").get_event_loop().run_until_complete(
            aggregator.aggregate_custom_with_previous(start, end)
        )
    )
    if not period_data:
        print("No data for the requested period.")
        return 1

    prs_service = PresentationService(settings)
    totals = prs_service._calculate_totals(period_data)
    
    # Get daily data for line charts
    daily_data = __import__("asyncio").get_event_loop().run_until_complete(
        aggregator.get_daily_series(start, end)
    )
    
    # Create PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    logo_path = getattr(settings, 'pptx_logo_path', '')
    
    # === 1. МОЩНЫЙ ТИТУЛЬНЫЙ ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Драматический градиент
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    add_gradient_fill(bg1, "#0D4F1A", "#2E7D32")
    bg1.line.fill.background()
    
    # Декоративные элементы с тенями
    circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.3), Inches(2.5), Inches(2.5))
    add_gradient_fill(circle1, "#4CAF50", "#66BB6A")
    circle1.line.fill.background()
    circle1.fill.transparency = 0.2
    add_shadow(circle1)
    
    diamond = slide1.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(10), Inches(5), Inches(2), Inches(2))
    add_gradient_fill(diamond, "#81C784", "#A5D6A7")
    diamond.line.fill.background()
    diamond.fill.transparency = 0.4
    add_shadow(diamond)
    
    # Логотип
    add_logo_to_slide(slide1, prs, logo_path)
    
    # Логотип справа сверху
    add_logo_to_slide(slide1, prs, logo_path)
    
    # Компактное название офиса сверху слева
    office_box = slide1.shapes.add_textbox(Inches(1), Inches(0.3), Inches(6), Inches(0.5))
    office_box.text_frame.text = getattr(settings, 'office_name', 'БАНКОВСКИЕ ГАРАНТИИ').upper()
    office_box.text_frame.paragraphs[0].font.size = Pt(16)
    office_box.text_frame.paragraphs[0].font.name = "Roboto"
    office_box.text_frame.paragraphs[0].font.bold = True
    office_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#E8F5E8")
    office_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # КРУПНЫЙ заголовок с тенью
    title_box = slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "ОТЧЁТ ПО ЭФФЕКТИВНОСТИ"
    title_frame.paragraphs[0].font.size = Pt(58)
    title_frame.paragraphs[0].font.name = "Roboto"
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_shadow(title_box)
    
    # Период НАД отчётом (заметный)
    period_box = slide1.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9.33), Inches(0.8))
    period_frame = period_box.text_frame
    period_frame.text = f"{period_name}"
    period_frame.paragraphs[0].font.size = Pt(20)
    period_frame.paragraphs[0].font.name = "Roboto"
    period_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#C8E6C9")
    period_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Даты под заголовком
    dates_box = slide1.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.33), Inches(0.6))
    dates_box.text_frame.text = f"{start.strftime('%d.%m.%Y')} — {end.strftime('%d.%m.%Y')}"
    dates_box.text_frame.paragraphs[0].font.size = Pt(28)
    dates_box.text_frame.paragraphs[0].font.name = "Roboto"
    dates_box.text_frame.paragraphs[0].font.bold = True
    dates_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    dates_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # === 2. ГРАФИКИ ЗВОНКОВ (по вашему эскизу) ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Градиентный фон
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    add_gradient_fill(bg2, "#F1F8E9", "#E8F5E8")
    bg2.line.fill.background()
    
    add_logo_to_slide(slide2, prs, logo_path)
    
    # Заголовок с градиентом
    header2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    add_gradient_fill(header2, "#1B5E20", "#2E7D32")
    header2.line.fill.background()
    
    header2_text = slide2.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11.33), Inches(0.5))
    header2_text.text_frame.text = "📈 ДИНАМИКА ЗВОНКОВ (НЕДЕЛЯ)"
    header2_text.text_frame.paragraphs[0].font.size = Pt(24)
    header2_text.text_frame.paragraphs[0].font.name = "Roboto"
    header2_text.text_frame.paragraphs[0].font.bold = True
    header2_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header2_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # График звонков (ваш эскиз: линейный план/факт)
    calls_chart_path = create_calls_chart(daily_data, "calls_dynamic.png")
    if os.path.exists(calls_chart_path):
        slide2.shapes.add_picture(calls_chart_path, Inches(1), Inches(1.5), width=Inches(6), height=Inches(3.5))
    
    # AI‑комментарий справа с градиентом
    ai_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.5), Inches(5), Inches(5.5))
    add_gradient_fill(ai_card, "#FFFFFF", "#F1F8E9")
    ai_card.line.color.rgb = hex_to_rgb("#C5E1A5")
    ai_card.line.width = Pt(2)
    add_shadow(ai_card)
    
    # Заголовок AI
    ai_title = slide2.shapes.add_textbox(Inches(8), Inches(1.8), Inches(4), Inches(0.4))
    ai_title.text_frame.text = "🤖 АНАЛИЗ КОМАНДЫ"
    ai_title.text_frame.paragraphs[0].font.size = Pt(16)
    ai_title.text_frame.paragraphs[0].font.name = "Roboto"
    ai_title.text_frame.paragraphs[0].font.bold = True
    ai_title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#2E7D32")
    ai_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # AI текст (генерируем)
    ai_comment = __import__("asyncio").get_event_loop().run_until_complete(
        YandexGPTService(settings).generate_team_comment(totals, period_name)
    )
    ai_text_box = slide2.shapes.add_textbox(Inches(8), Inches(2.4), Inches(4), Inches(4))
    ai_text_box.text_frame.text = ai_comment[:220] + ("..." if len(ai_comment) > 220 else "")
    ai_text_box.text_frame.word_wrap = True
    for p in ai_text_box.text_frame.paragraphs:
        p.font.name = "Roboto"
        p.font.size = Pt(13)
        p.font.color.rgb = hex_to_rgb("#424242")
        p.line_spacing = 1.3
    
    # === 3. ТАБЛИЦА МЕНЕДЖЕРОВ 3×3 (ваш эскиз) ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    add_gradient_fill(bg3, "#FAFAFA", "#F0F4C3")
    bg3.line.fill.background()
    
    add_logo_to_slide(slide3, prs, logo_path)
    
    # Заголовок
    header3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    add_gradient_fill(header3, "#689F38", "#8BC34A")
    header3.line.fill.background()
    
    header3_text = slide3.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11.33), Inches(0.5))
    header3_text.text_frame.text = "👥 РЕЗУЛЬТАТЫ ПО МЕНЕДЖЕРАМ"
    header3_text.text_frame.paragraphs[0].font.size = Pt(24)
    header3_text.text_frame.paragraphs[0].font.name = "Roboto"
    header3_text.text_frame.paragraphs[0].font.bold = True
    header3_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header3_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Сетка менеджеров 3×3 с карточками
    managers = list(period_data.values())
    card_w, card_h = Inches(4), Inches(1.5)
    start_x, start_y = Inches(0.5), Inches(1.5)
    
    for i, mgr in enumerate(managers[:9]):  # Max 9 for 3x3
        col, row = i % 3, i // 3
        x = start_x + col * Inches(4.2)
        y = start_y + row * Inches(1.8)
        
        # Карточка с градиентом
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        
        # Цвет карточки по эффективности
        efficiency = (mgr.calls_percentage + mgr.leads_volume_percentage) / 2
        if efficiency >= 80:
            add_gradient_fill(card, "#4CAF50", "#66BB6A")  # Зелёный
        elif efficiency >= 60:
            add_gradient_fill(card, "#FF9800", "#FFB74D")  # Жёлтый
        else:
            add_gradient_fill(card, "#F44336", "#EF5350")  # Красный
            
        card.line.fill.background()
        add_shadow(card)
        
        # Имя менеджера
        name_box = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.4), Inches(0.4))
        name_box.text_frame.text = mgr.name
        name_box.text_frame.paragraphs[0].font.size = Pt(14)
        name_box.text_frame.paragraphs[0].font.name = "Roboto"
        name_box.text_frame.paragraphs[0].font.bold = True
        name_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        name_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Показатели
        stats_text = f"📞 {mgr.calls_fact}/{mgr.calls_plan} • 📝 {mgr.leads_units_fact} • 💰 {mgr.leads_volume_fact:.1f}".replace(".", ",")
        stats_box = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), card_w - Inches(0.4), Inches(0.6))
        stats_box.text_frame.text = stats_text
        stats_box.text_frame.paragraphs[0].font.size = Pt(11)
        stats_box.text_frame.paragraphs[0].font.name = "Roboto"
        stats_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        stats_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # === 4. SPIDER ДИАГРАММА (ваш эскиз) ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    add_gradient_fill(bg4, "#F3E5F5", "#E1BEE7")
    bg4.line.fill.background()
    
    add_logo_to_slide(slide4, prs, logo_path)
    
    # Заголовок
    header4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    add_gradient_fill(header4, "#7B1FA2", "#9C27B0")
    header4.line.fill.background()
    
    header4_text = slide4.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11.33), Inches(0.5))
    header4_text.text_frame.text = "🕷️ ПРОФИЛЬ ЭФФЕКТИВНОСТИ"
    header4_text.text_frame.paragraphs[0].font.size = Pt(24)
    header4_text.text_frame.paragraphs[0].font.name = "Roboto"
    header4_text.text_frame.paragraphs[0].font.bold = True
    header4_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header4_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Spider диаграмма (первый менеджер vs среднее)
    if managers:
        first_mgr = managers[0]
        cnt = max(len(managers), 1)
        avg_mgr_data = {
            'calls_fact': totals['calls_fact'] / cnt,
            'new_calls': totals['new_calls'] / cnt,
            'leads_units_fact': totals['leads_units_fact'] / cnt,
            'leads_volume_fact': totals['leads_volume_fact'] / cnt,
            'approved_volume': totals['approved_volume'] / cnt,
            'issued_volume': totals['issued_volume'] / cnt,
        }
        
        spider_path = create_spider_chart(first_mgr, avg_mgr_data, first_mgr.name, "spider_analysis.png")
        if os.path.exists(spider_path):
            slide4.shapes.add_picture(spider_path, Inches(2), Inches(1.5), width=Inches(9), height=Inches(5))
    
    # === 5. СТОЛБЧАТАЯ ДИАГРАММА МЕНЕДЖЕРОВ ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    add_gradient_fill(bg5, "#E3F2FD", "#BBDEFB")
    bg5.line.fill.background()
    
    add_logo_to_slide(slide5, prs, logo_path)
    
    # Заголовок
    header5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    add_gradient_fill(header5, "#1976D2", "#2196F3")
    header5.line.fill.background()
    
    header5_text = slide5.shapes.add_textbox(Inches(1), Inches(0.25), Inches(11.33), Inches(0.5))
    header5_text.text_frame.text = "📊 СРАВНЕНИЕ КОМАНДЫ"
    header5_text.text_frame.paragraphs[0].font.size = Pt(24)
    header5_text.text_frame.paragraphs[0].font.name = "Roboto"
    header5_text.text_frame.paragraphs[0].font.bold = True
    header5_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header5_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Столбчатая диаграмма
    bar_path = create_bar_chart(managers, "managers_comparison.png")
    if os.path.exists(bar_path):
        slide5.shapes.add_picture(bar_path, Inches(1.5), Inches(1.5), width=Inches(10), height=Inches(5))
    
    # Save
    out = f"powerhouse_presentation_{start.strftime('%Y%m%d')}_{end.strftime('%Y%m%d')}.pptx"
    prs.save(out)
    print(f"🚀 POWERHOUSE PPTX saved: {out}")
    print("Includes: dramatic gradients, shadows, line charts, spider, 3x3 manager grid, bar charts")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

