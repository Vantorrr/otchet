#!/usr/bin/env python3
"""Генератор презентации СТРОГО по брифу клиента (9 слайдов, точная палитра, таблицы, диаграммы)."""
from __future__ import annotations

import os
import sys
from datetime import datetime

from dotenv import load_dotenv

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from bot.config import Settings
from bot.services.di import Container
from bot.services.data_aggregator import DataAggregatorService
from bot.services.presentation import PresentationService
from bot.services.yandex_gpt import YandexGPTService

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import plotly.graph_objects as go


# === УТИЛИТЫ ===
def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)

PRIMARY = "#2E7D32"
ALERT = "#C62828"
ACCENT2 = "#FF8A65"
TEXT_MAIN = "#222222"
TEXT_MUTED = "#6B6B6B"
CARD_BG = "#F5F5F5"
SLIDE_BG = "#FFFFFF"

def add_logo(slide, prs, logo_path):
    """Добавить логотип справа сверху."""
    # Check multiple possible paths
    possible_paths = [
        logo_path,
        os.path.join(os.path.dirname(__file__), "..", "Логотип.png"),
        "Логотип.png",
        "assets/logo.png"
    ]
    
    for path in possible_paths:
        if path and os.path.exists(path):
            try:
                slide.shapes.add_picture(
                    path, 
                    prs.slide_width - Inches(2.2), 
                    Inches(0.2), 
                    height=Inches(0.75)
                )
                return
            except Exception as e:
                continue

def create_donut_chart(totals, path="donut.png"):
    labels = ['Повторные звонки', 'Заявки шт', 'Заявки млн', 'Выдано']
    values = [totals['calls_fact'], totals['leads_units_fact'], totals['leads_volume_fact']*10, totals['issued_volume']*10]
    colors = [PRIMARY, ACCENT2, '#81C784', '#AED581']
    fig = go.Figure(data=[go.Pie(labels=labels, values=values, hole=.5, marker_colors=colors)])
    fig.update_layout(
        title="Распределение активности", title_font_size=14, title_font_family="Roboto",
        font=dict(family="Roboto", size=11), showlegend=True, width=400, height=350,
        paper_bgcolor='rgba(0,0,0,0)'
    )
    fig.write_image(path)
    return path

def create_comparison_bars(prev, cur, path="compare_bars.png"):
    categories = ['Звонки', 'Заявки шт', 'Заявки млн']
    prev_vals = [prev['calls_fact'], prev['leads_units_fact'], prev['leads_volume_fact']]
    cur_vals = [cur['calls_fact'], cur['leads_units_fact'], cur['leads_volume_fact']]
    
    fig = go.Figure()
    fig.add_trace(go.Bar(name='Предыдущий', x=categories, y=prev_vals, marker_color=ACCENT2))
    fig.add_trace(go.Bar(name='Текущий', x=categories, y=cur_vals, marker_color=PRIMARY))
    fig.update_layout(barmode='group', title="Сравнение периодов", title_font_family="Roboto",
                      font=dict(family="Roboto", size=11), width=700, height=350, paper_bgcolor='rgba(0,0,0,0)')
    fig.write_image(path)
    return path

def create_line_dynamics(daily_data, path="line_dynamics.png"):
    dates = [d['date'] for d in daily_data]
    plan = [d['leads_volume_plan'] for d in daily_data]
    fact = [d['leads_volume_fact'] for d in daily_data]
    issued = [d['issued_volume'] for d in daily_data]
    
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=dates, y=plan, mode='lines+markers', name='План', line=dict(color=PRIMARY, width=2)))
    fig.add_trace(go.Scatter(x=dates, y=fact, mode='lines+markers', name='Факт', line=dict(color=ACCENT2, width=2)))
    fig.add_trace(go.Scatter(x=dates, y=issued, mode='lines+markers', name='Выдано', line=dict(color=ALERT, width=2)))
    
    fig.update_layout(title="Динамика (млн)", title_font_family="Roboto", xaxis_title="Дни", yaxis_title="Млн ₽",
                      font=dict(family="Roboto", size=11), width=700, height=350, paper_bgcolor='rgba(0,0,0,0)')
    fig.write_image(path)
    return path


def main():
    load_dotenv()
    os.environ.setdefault("GOOGLE_APPLICATION_CREDENTIALS", "service_account.json")
    os.environ.setdefault("BOT_TOKEN", "dummy")
    os.environ.setdefault("OFFICE_NAME", "Банковские гарантии")

    settings = Settings.load()
    Container.init(settings)

    start = datetime.strptime("2025-09-01", "%Y-%m-%d").date()
    end = datetime.strptime("2025-09-07", "%Y-%m-%d").date()
    print("🔨 Генерация презентации СТРОГО ПО БРИФУ (9 слайдов)...")

    aggregator = DataAggregatorService(Container.get().sheets)
    period_data, prev_data, period_name, s1, e1, s2, e2 = __import__("asyncio").get_event_loop().run_until_complete(
        aggregator.aggregate_custom_with_previous(start, end)
    )
    if not period_data:
        print("Нет данных.")
        return 1

    prs_service = PresentationService(settings)
    totals = prs_service._calculate_totals(period_data)
    prev_totals = prs_service._calculate_totals(prev_data) if prev_data else {}
    
    daily = __import__("asyncio").get_event_loop().run_until_complete(aggregator.get_daily_series(start, end))
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    logo = getattr(settings, 'pptx_logo_path', '')
    margin = Cm(1.5)
    
    # === 1. ТИТУЛЬНЫЙ ===
    print("  Слайд 1/9: Титульный...")
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg1.line.fill.background()
    
    add_logo(s1, prs, logo)
    
    # Название офиса сверху
    office = s1.shapes.add_textbox(margin, Inches(0.5), Inches(8), Inches(0.5))
    office.text_frame.text = getattr(settings, 'office_name', 'Банковские гарантии').upper()
    office.text_frame.paragraphs[0].font.name = "Roboto"
    office.text_frame.paragraphs[0].font.size = Pt(16)
    office.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(TEXT_MUTED)
    
    # Заголовок
    title = s1.shapes.add_textbox(margin, Inches(2.5), prs.slide_width - 2*margin, Inches(1.5))
    title.text_frame.text = "ОТЧЁТ ПО ПРОДАЖАМ"
    title.text_frame.paragraphs[0].font.name = "Roboto"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Период
    period = s1.shapes.add_textbox(margin, Inches(4.2), prs.slide_width - 2*margin, Inches(0.7))
    period.text_frame.text = f"{period_name}"
    period.text_frame.paragraphs[0].font.name = "Roboto"
    period.text_frame.paragraphs[0].font.size = Pt(20)
    period.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(TEXT_MAIN)
    period.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Даты
    dates = s1.shapes.add_textbox(margin, Inches(5), prs.slide_width - 2*margin, Inches(0.5))
    dates.text_frame.text = f"{start.strftime('%d.%m.%Y')} — {end.strftime('%d.%m.%Y')}"
    dates.text_frame.paragraphs[0].font.name = "Roboto"
    dates.text_frame.paragraphs[0].font.size = Pt(16)
    dates.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(TEXT_MUTED)
    dates.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # === 2. СВОДКА (KEY METRICS) ===
    print("  Слайд 2/9: Сводка команды...")
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg2.line.fill.background()
    add_logo(s2, prs, logo)
    
    # Заголовок
    h2 = s2.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h2.text_frame.text = "Ключевые показатели команды"
    h2.text_frame.paragraphs[0].font.name = "Roboto"
    h2.text_frame.paragraphs[0].font.size = Pt(28)
    h2.text_frame.paragraphs[0].font.bold = True
    h2.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Таблица с зеброй
    rows, cols = 7, 4
    tbl = s2.shapes.add_table(rows, cols, margin, Inches(1.3), Inches(11.33), Inches(4.5)).table
    
    # Заголовки
    headers = ["Показатель", "План", "Факт", "Конв (%)"]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
        for p in cell.text_frame.paragraphs:
            p.font.name = "Roboto"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    
    # Данные
    data_rows = [
        ("Повторные звонки", f"{int(totals['calls_plan']):,}".replace(",", " "), f"{int(totals['calls_fact']):,}".replace(",", " "), f"{totals['calls_percentage']:.1f}%".replace(".", ",")),
        ("Заявки, шт", f"{int(totals['leads_units_plan']):,}".replace(",", " "), f"{int(totals['leads_units_fact']):,}".replace(",", " "), f"{totals['leads_units_percentage']:.1f}%".replace(".", ",")),
        ("Заявки, млн", f"{totals['leads_volume_plan']:.1f}".replace(".", ","), f"{totals['leads_volume_fact']:.1f}".replace(".", ","), f"{totals['leads_volume_percentage']:.1f}%".replace(".", ",")),
        ("Одобрено, млн", "—", f"{totals['approved_volume']:.1f}".replace(".", ","), "—"),
        ("Выдано, млн", "—", f"{totals['issued_volume']:.1f}".replace(".", ","), "—"),
        ("Новые звонки", "—", f"{int(totals['new_calls']):,}".replace(",", " "), "—"),
    ]
    
    for r, (name, plan, fact, conv) in enumerate(data_rows, start=1):
        for c, val in enumerate([name, plan, fact, conv]):
            cell = tbl.cell(r, c)
            cell.text = val
            # Зебра
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(12)
                p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    
    # Donut НИЖЕ таблицы (отдельно)
    donut_path = create_donut_chart(totals, "donut_metrics.png")
    if os.path.exists(donut_path):
        s2.shapes.add_picture(donut_path, Inches(4), Inches(6.2), width=Inches(5), height=Inches(1.2))
    
    # === 3. AI‑КОММЕНТАРИЙ КОМАНДЫ ===
    print("  Слайд 3/9: AI-комментарий по команде...")
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg3.line.fill.background()
    add_logo(s3, prs, logo)
    
    h3 = s3.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h3.text_frame.text = "Анализ и рекомендации"
    h3.text_frame.paragraphs[0].font.name = "Roboto"
    h3.text_frame.paragraphs[0].font.size = Pt(28)
    h3.text_frame.paragraphs[0].font.bold = True
    h3.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # AI текст
    ai_comment = __import__("asyncio").get_event_loop().run_until_complete(
        YandexGPTService(settings).generate_team_comment(totals, period_name)
    )
    ai_box = s3.shapes.add_textbox(margin, Inches(1.5), prs.slide_width - 2*margin, Inches(5))
    ai_box.text_frame.text = ai_comment
    ai_box.text_frame.word_wrap = True
    for p in ai_box.text_frame.paragraphs:
        p.font.name = "Roboto"
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
        p.line_spacing = 1.3
        p.alignment = PP_ALIGN.LEFT
    
    # === 4. СРАВНЕНИЕ С ПРЕДЫДУЩИМ (3 ГРАФИКА) ===
    print("  Слайд 4/9: Сравнение периодов...")
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg4 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg4.line.fill.background()
    add_logo(s4, prs, logo)
    
    h4 = s4.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h4.text_frame.text = "Сравнение с предыдущим периодом"
    h4.text_frame.paragraphs[0].font.name = "Roboto"
    h4.text_frame.paragraphs[0].font.size = Pt(28)
    h4.text_frame.paragraphs[0].font.bold = True
    h4.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if prev_totals:
        compare_path = create_comparison_bars(prev_totals, totals, "comparison_bars.png")
        if os.path.exists(compare_path):
            s4.shapes.add_picture(compare_path, Inches(1), Inches(1.5), width=Inches(5.5), height=Inches(3))
    
    line_path = create_line_dynamics(daily, "dynamics_line.png")
    if os.path.exists(line_path):
        s4.shapes.add_picture(line_path, Inches(7), Inches(1.5), width=Inches(5.5), height=Inches(3))
    
    # Donut changes
    if os.path.exists("donut_metrics.png"):
        s4.shapes.add_picture("donut_metrics.png", Inches(4), Inches(5), width=Inches(5), height=Inches(2))
    
    # === 5. ТОП-2 И АНТИТОП-2 ===
    print("  Слайд 5/9: Рейтинг...")
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg5 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg5.line.fill.background()
    add_logo(s5, prs, logo)
    
    h5 = s5.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h5.text_frame.text = "Рейтинг эффективности"
    h5.text_frame.paragraphs[0].font.name = "Roboto"
    h5.text_frame.paragraphs[0].font.size = Pt(28)
    h5.text_frame.paragraphs[0].font.bold = True
    h5.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Ranking
    scored = []
    for m in period_data.values():
        cp = (m.calls_fact/m.calls_plan*100) if m.calls_plan else 0
        vp = (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0
        scored.append((0.5*cp+0.5*vp, m.name))
    scored.sort(reverse=True)
    best = [n for _, n in scored[:2]]
    worst = [n for _, n in list(reversed(scored[-2:]))]
    
    # ТОП-2
    top_card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin, Inches(1.5), Inches(6), Inches(5))
    top_card.fill.solid()
    top_card.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
    top_card.line.fill.background()
    
    top_t = s5.shapes.add_textbox(margin + Inches(0.5), Inches(2), Inches(5), Inches(0.5))
    top_t.text_frame.text = "🏆 ЛИДЕРЫ ПЕРИОДА"
    top_t.text_frame.paragraphs[0].font.name = "Roboto"
    top_t.text_frame.paragraphs[0].font.size = Pt(22)
    top_t.text_frame.paragraphs[0].font.bold = True
    top_t.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    for i, name in enumerate(best[:2]):
        nb = s5.shapes.add_textbox(margin + Inches(0.5), Inches(2.8 + i*1), Inches(5), Inches(0.8))
        nb.text_frame.text = f"{i+1}. {name}"
        nb.text_frame.paragraphs[0].font.name = "Roboto"
        nb.text_frame.paragraphs[0].font.size = Pt(18)
        nb.text_frame.paragraphs[0].font.bold = True
        nb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # АНТИТОП-2
    bot_card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(6), Inches(5))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = hex_to_rgb(ALERT)
    bot_card.line.fill.background()
    
    bot_t = s5.shapes.add_textbox(Inches(7.5), Inches(2), Inches(5), Inches(0.5))
    bot_t.text_frame.text = "⚠️ ТРЕБУЮТ ВНИМАНИЯ"
    bot_t.text_frame.paragraphs[0].font.name = "Roboto"
    bot_t.text_frame.paragraphs[0].font.size = Pt(22)
    bot_t.text_frame.paragraphs[0].font.bold = True
    bot_t.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    for i, name in enumerate(worst[:2]):
        nb = s5.shapes.add_textbox(Inches(7.5), Inches(2.8 + i*1), Inches(5), Inches(0.8))
        nb.text_frame.text = f"{i+1}. {name}"
        nb.text_frame.paragraphs[0].font.name = "Roboto"
        nb.text_frame.paragraphs[0].font.size = Pt(18)
        nb.text_frame.paragraphs[0].font.bold = True
        nb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # === 6. ОБЩАЯ ТАБЛИЦА ВСЕХ МЕНЕДЖЕРОВ ===
    print("  Слайд 6/9: Таблица всех менеджеров...")
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    bg6 = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg6.line.fill.background()
    add_logo(s6, prs, logo)
    
    h6 = s6.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h6.text_frame.text = "Результаты по менеджерам"
    h6.text_frame.paragraphs[0].font.name = "Roboto"
    h6.text_frame.paragraphs[0].font.size = Pt(28)
    h6.text_frame.paragraphs[0].font.bold = True
    h6.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h6.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Таблица менеджеров
    mgr_rows = len(period_data) + 1
    tbl6 = s6.shapes.add_table(mgr_rows, 5, margin, Inches(1.3), Inches(11.33), Inches(5.5)).table
    
    mgr_headers = ["Менеджер", "Звонки", "Заявки шт", "Заявки млн", "Выдано млн"]
    for c, h in enumerate(mgr_headers):
        cell = tbl6.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
        for p in cell.text_frame.paragraphs:
            p.font.name = "Roboto"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    
    for r, mgr in enumerate(period_data.values(), start=1):
        row_data = [
            mgr.name,
            f"{mgr.calls_fact}/{mgr.calls_plan}",
            f"{int(mgr.leads_units_fact):,}".replace(",", " "),
            f"{mgr.leads_volume_fact:.1f}".replace(".", ","),
            f"{mgr.issued_volume:.1f}".replace(".", ",")
        ]
        for c, val in enumerate(row_data):
            cell = tbl6.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(12)
                p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    
    # === 7. ИНДИВИДУАЛЬНЫЕ КАРТОЧКИ МЕНЕДЖЕРОВ ===
    print("  Слайд 7/9: Карточки менеджеров...")
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    bg7 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg7.line.fill.background()
    add_logo(s7, prs, logo)
    
    h7 = s7.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h7.text_frame.text = "Индивидуальные показатели"
    h7.text_frame.paragraphs[0].font.name = "Roboto"
    h7.text_frame.paragraphs[0].font.size = Pt(28)
    h7.text_frame.paragraphs[0].font.bold = True
    h7.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h7.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Сетка карточек менеджеров 3×3
    managers = list(period_data.values())
    for i, mgr in enumerate(managers[:9]):
        col, row = i % 3, i // 3
        x = margin + col * Inches(4.2)
        y = Inches(1.5) + row * Inches(1.9)
        
        # Карточка с цветом по эффективности
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(1.7))
        eff = (mgr.calls_percentage + mgr.leads_volume_percentage) / 2
        if eff >= 80:
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
        elif eff >= 60:
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb("#FF9800")
        else:
            card.fill.solid()
            card.fill.fore_color.rgb = hex_to_rgb(ALERT)
        card.line.fill.background()
        
        # Имя
        name_box = s7.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.5))
        name_box.text_frame.text = mgr.name
        name_box.text_frame.paragraphs[0].font.name = "Roboto"
        name_box.text_frame.paragraphs[0].font.size = Pt(14)
        name_box.text_frame.paragraphs[0].font.bold = True
        name_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        name_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Показатели
        stats = f"📞 {mgr.calls_fact} • 📝 {mgr.leads_units_fact} • 💰 {mgr.leads_volume_fact:.1f}".replace(".", ",")
        stats_box = s7.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), Inches(3.5), Inches(0.7))
        stats_box.text_frame.text = stats
        stats_box.text_frame.paragraphs[0].font.name = "Roboto"
        stats_box.text_frame.paragraphs[0].font.size = Pt(12)
        stats_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        stats_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # === 8. ДИНАМИКА КЛЮЧЕВЫХ МЕТРИК ===
    print("  Слайд 8/9: Динамика метрик...")
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    bg8 = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg8.line.fill.background()
    add_logo(s8, prs, logo)
    
    h8 = s8.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h8.text_frame.text = "Динамика ключевых метрик"
    h8.text_frame.paragraphs[0].font.name = "Roboto"
    h8.text_frame.paragraphs[0].font.size = Pt(28)
    h8.text_frame.paragraphs[0].font.bold = True
    h8.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h8.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Линейный график динамики
    if os.path.exists(line_path):
        s8.shapes.add_picture(line_path, Inches(2), Inches(1.5), width=Inches(9), height=Inches(5))
    
    # === 9. ВЫВОДЫ И РЕКОМЕНДАЦИИ AI ===
    print("  Слайд 9/9: Итоги и рекомендации...")
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    bg9 = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = hex_to_rgb(SLIDE_BG)
    bg9.line.fill.background()
    add_logo(s9, prs, logo)
    
    h9 = s9.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
    h9.text_frame.text = "Выводы и рекомендации"
    h9.text_frame.paragraphs[0].font.name = "Roboto"
    h9.text_frame.paragraphs[0].font.size = Pt(28)
    h9.text_frame.paragraphs[0].font.bold = True
    h9.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    h9.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # AI итоговые рекомендации
    final_ai = __import__("asyncio").get_event_loop().run_until_complete(
        YandexGPTService(settings).generate_team_comment(totals, f"Итоги {period_name}")
    )
    final_box = s9.shapes.add_textbox(margin, Inches(1.5), prs.slide_width - 2*margin, Inches(5))
    final_box.text_frame.text = f"🎯 КЛЮЧЕВЫЕ ВЫВОДЫ:\n\n{final_ai}\n\n📌 СЛЕДУЮЩИЕ ШАГИ:\n• Усилить работу с отстающими\n• Масштабировать успешные практики\n• Оптимизировать процессы"
    final_box.text_frame.word_wrap = True
    for p in final_box.text_frame.paragraphs:
        p.font.name = "Roboto"
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
        p.line_spacing = 1.3
    
    # Сохранение
    out = f"brief_compliant_{start.strftime('%Y%m%d')}_{end.strftime('%Y%m%d')}.pptx"
    prs.save(out)
    print(f"\n✅ ПРЕЗЕНТАЦИЯ ПО БРИФУ: {out}")
    print("Включено: 9 слайдов, точная палитра, зебра-таблицы, donut/bar/line, AI-комментарии, логотип на всех")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
