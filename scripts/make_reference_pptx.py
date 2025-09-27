#!/usr/bin/env python3
"""Generate reference-style PPTX locally to show final design."""
from __future__ import annotations

import os
import sys
import io
from datetime import datetime, date

from dotenv import load_dotenv

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from bot.config import Settings
from bot.services.di import Container
from bot.services.data_aggregator import DataAggregatorService
from bot.services.presentation import PresentationService

# PPTX imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def main():
    load_dotenv()
    os.environ.setdefault("GOOGLE_APPLICATION_CREDENTIALS", "service_account.json")
    os.environ.setdefault("BOT_TOKEN", "dummy-token-for-offline")

    settings = Settings.load()
    Container.init(settings)

    start = datetime.strptime("2025-09-01", "%Y-%m-%d").date()
    end = datetime.strptime("2025-09-07", "%Y-%m-%d").date()
    print(f"Building REFERENCE-STYLE PPTX for {start}..{end}")

    aggregator = DataAggregatorService(Container.get().sheets)
    period_data, prev_data, period_name, start_date, end_date, prev_start, prev_end = (
        __import__("asyncio").get_event_loop().run_until_complete(
            aggregator.aggregate_custom_with_previous(start, end)
        )
    )
    if not period_data:
        print("No data for the requested period.")
        return 1

    prs_service = PresentationService(settings)
    totals = prs_service._calculate_totals(period_data)
    
    # Create PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # === TITLE SLIDE ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Emerald background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb("#2E7D32")
    bg_shape.line.fill.background()
    
    # Geometric accents
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(1.2), Inches(1.2))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = hex_to_rgb("#4CAF50")
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(5.5), Inches(1.5), Inches(1.5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    circle2.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.33), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = f"{getattr(settings, 'office_name', 'БАНКОВСКИЕ ГАРАНТИИ')}\n\nОТЧЕТ ПО ПРОДАЖАМ"
    for p in title_frame.paragraphs:
        p.font.size = Pt(42)
        p.font.name = "Roboto"
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # Period subtitle
    period_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.33), Inches(1.5))
    period_frame = period_box.text_frame
    period_frame.text = f"{period_name}\n{start.strftime('%d.%m.%Y')} — {end.strftime('%d.%m.%Y')}"
    for p in period_frame.paragraphs:
        p.font.size = Pt(24)
        p.font.name = "Roboto"
        p.font.color.rgb = hex_to_rgb("#F5F5DC")
        p.alignment = PP_ALIGN.CENTER
    
    # === METRICS DASHBOARD ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Light background
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = hex_to_rgb("#F5F5DC")
    bg2.line.fill.background()
    
    # Header bar
    header_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = hex_to_rgb("#2E7D32")
    header_bar.line.fill.background()
    
    # Header text
    header_text = slide2.shapes.add_textbox(Inches(1), Inches(0.2), Inches(11.33), Inches(0.6))
    header_text.text_frame.text = "КЛЮЧЕВЫЕ ПОКАЗАТЕЛИ"
    header_text.text_frame.paragraphs[0].font.size = Pt(24)
    header_text.text_frame.paragraphs[0].font.name = "Roboto"
    header_text.text_frame.paragraphs[0].font.bold = True
    header_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Metrics cards (3x2 grid)
    metrics = [
        ("ПОВТОРНЫЕ\nЗВОНКИ", f"{int(totals.get('calls_fact', 0)):,}".replace(",", " "), f"{totals.get('calls_percentage', 0):.1f}%".replace(".", ",")),
        ("ЗАЯВКИ\n(ШТ)", f"{int(totals.get('leads_units_fact', 0)):,}".replace(",", " "), f"{totals.get('leads_units_percentage', 0):.1f}%".replace(".", ",")),
        ("ЗАЯВКИ\n(МЛН)", f"{totals.get('leads_volume_fact', 0):.1f}".replace(".", ","), f"{totals.get('leads_volume_percentage', 0):.1f}%".replace(".", ",")),
        ("ОДОБРЕНО\n(МЛН)", f"{totals.get('approved_volume', 0):.1f}".replace(".", ","), "—"),
        ("ВЫДАНО\n(МЛН)", f"{totals.get('issued_volume', 0):.1f}".replace(".", ","), "—"),
        ("НОВЫЕ\nЗВОНКИ", f"{int(totals.get('new_calls', 0)):,}".replace(",", " "), "—"),
    ]
    
    card_w = Inches(3.8)
    card_h = Inches(1.5)
    for i, (label, value, pct) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = Inches(1 + col * 4.1)
        y = Inches(1.3 + row * 1.8)
        
        # Card background
        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = hex_to_rgb("#E0E0E0")
        
        # Accent top bar
        accent = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = hex_to_rgb("#4CAF50")
        accent.line.fill.background()
        
        # Label
        label_box = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.4))
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(11)
        label_box.text_frame.paragraphs[0].font.name = "Roboto"
        label_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#757575")
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Value
        value_box = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.5))
        value_box.text_frame.text = value
        value_box.text_frame.paragraphs[0].font.size = Pt(22)
        value_box.text_frame.paragraphs[0].font.name = "Roboto"
        value_box.text_frame.paragraphs[0].font.bold = True
        value_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#2E7D32")
        value_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Percentage
        if pct != "—":
            pct_box = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(1.2), card_w - Inches(0.4), Inches(0.25))
            pct_box.text_frame.text = pct
            pct_box.text_frame.paragraphs[0].font.size = Pt(14)
            pct_box.text_frame.paragraphs[0].font.name = "Roboto"
            pct_box.text_frame.paragraphs[0].font.bold = True
            pct_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#4CAF50")
            pct_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    # === RANKING SLIDE ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Light background
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = hex_to_rgb("#F5F5DC")
    bg3.line.fill.background()
    
    # Header
    header3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header3.fill.solid()
    header3.fill.fore_color.rgb = hex_to_rgb("#2E7D32")
    header3.line.fill.background()
    
    header3_text = slide3.shapes.add_textbox(Inches(1), Inches(0.2), Inches(11.33), Inches(0.6))
    header3_text.text_frame.text = "РЕЙТИНГ ЭФФЕКТИВНОСТИ"
    header3_text.text_frame.paragraphs[0].font.size = Pt(24)
    header3_text.text_frame.paragraphs[0].font.name = "Roboto"
    header3_text.text_frame.paragraphs[0].font.bold = True
    header3_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header3_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Build ranking
    scored = []
    for m in period_data.values():
        calls_pct = (m.calls_fact/m.calls_plan*100) if m.calls_plan else 0
        vol_pct = (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0
        scored.append((0.5*calls_pct+0.5*vol_pct, m.name))
    scored.sort(reverse=True)
    best = [n for _, n in scored[:2]]
    worst = [n for _, n in list(reversed(scored[-2:]))]
    
    # Best performers card (large emerald)
    best_card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(5.5), Inches(4))
    best_card.fill.solid()
    best_card.fill.fore_color.rgb = hex_to_rgb("#2E7D32")
    best_card.line.fill.background()
    
    # Crown icon background
    crown_bg = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), Inches(1.8), Inches(0.8), Inches(0.8))
    crown_bg.fill.solid()
    crown_bg.fill.fore_color.rgb = hex_to_rgb("#4CAF50")
    crown_bg.line.fill.background()
    
    # Best title
    best_title = slide3.shapes.add_textbox(Inches(2.5), Inches(1.9), Inches(3.7), Inches(0.6))
    best_title.text_frame.text = "🏆 ЛИДЕРЫ ПЕРИОДА"
    best_title.text_frame.paragraphs[0].font.size = Pt(20)
    best_title.text_frame.paragraphs[0].font.name = "Roboto"
    best_title.text_frame.paragraphs[0].font.bold = True
    best_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    best_title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Best names
    for i, name in enumerate(best[:2]):
        name_box = slide3.shapes.add_textbox(Inches(1.3), Inches(2.8 + i*0.8), Inches(4.9), Inches(0.6))
        name_box.text_frame.text = f"{i+1}. {name}"
        name_box.text_frame.paragraphs[0].font.size = Pt(18)
        name_box.text_frame.paragraphs[0].font.name = "Roboto"
        name_box.text_frame.paragraphs[0].font.bold = True
        name_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        desc_box = slide3.shapes.add_textbox(Inches(1.3), Inches(3.2 + i*0.8), Inches(4.9), Inches(0.4))
        desc_box.text_frame.text = "высокая результативность"
        desc_box.text_frame.paragraphs[0].font.size = Pt(12)
        desc_box.text_frame.paragraphs[0].font.name = "Roboto"
        desc_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#F5F5DC")
    
    # Worst performers card (red-orange)
    worst_card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), Inches(5.5), Inches(4))
    worst_card.fill.solid()
    worst_card.fill.fore_color.rgb = hex_to_rgb("#D32F2F")
    worst_card.line.fill.background()
    
    # Warning icon
    warn_bg = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), Inches(1.8), Inches(0.8), Inches(0.8))
    warn_bg.fill.solid()
    warn_bg.fill.fore_color.rgb = hex_to_rgb("#FF5722")
    warn_bg.line.fill.background()
    
    # Worst title
    worst_title = slide3.shapes.add_textbox(Inches(8.5), Inches(1.9), Inches(3.7), Inches(0.6))
    worst_title.text_frame.text = "⚠️ ТРЕБУЮТ ВНИМАНИЯ"
    worst_title.text_frame.paragraphs[0].font.size = Pt(20)
    worst_title.text_frame.paragraphs[0].font.name = "Roboto"
    worst_title.text_frame.paragraphs[0].font.bold = True
    worst_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    worst_title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Worst names
    for i, name in enumerate(worst[:2]):
        name_box = slide3.shapes.add_textbox(Inches(7.3), Inches(2.8 + i*0.8), Inches(4.9), Inches(0.6))
        name_box.text_frame.text = f"{i+1}. {name}"
        name_box.text_frame.paragraphs[0].font.size = Pt(18)
        name_box.text_frame.paragraphs[0].font.name = "Roboto"
        name_box.text_frame.paragraphs[0].font.bold = True
        name_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        desc_box = slide3.shapes.add_textbox(Inches(7.3), Inches(3.2 + i*0.8), Inches(4.9), Inches(0.4))
        desc_box.text_frame.text = "снижение показателей"
        desc_box.text_frame.paragraphs[0].font.size = Pt(12)
        desc_box.text_frame.paragraphs[0].font.name = "Roboto"
        desc_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb("#FFE0B2")
    
    # Save
    out = f"reference_business_{start.strftime('%Y%m%d')}_{end.strftime('%Y%m%d')}.pptx"
    prs.save(out)
    print(f"REFERENCE PPTX saved: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
