"""Premium PPTX presentation generator - 9 slides with charts, diagrams, AI analysis."""
from __future__ import annotations

import os
import io
import asyncio
from datetime import datetime, date
from typing import Dict, Any, List, Tuple, Optional
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
import plotly.graph_objects as go

from bot.services.yandex_gpt import YandexGPTService
from bot.services.presentation import ManagerData
from bot.config import Settings


# Брендинг
PRIMARY = "#2E7D32"
ALERT = "#C62828"
ACCENT2 = "#FF8A65"
TEXT_MAIN = "#222222"
TEXT_MUTED = "#6B6B6B"
CARD_BG = "#F5F5F5"
SLIDE_BG = "#FFFFFF"


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def add_gradient_bg(slide, prs, color_theme="default"):
    """Add gradient background - default white, or themed (green/purple/blue)."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fill = bg.fill
    
    if color_theme == "green":
        fill.gradient()
        fill.gradient_angle = 90.0
        fill.gradient_stops[0].color.rgb = hex_to_rgb(PRIMARY)
        fill.gradient_stops[1].color.rgb = hex_to_rgb("#1B5E20")
    elif color_theme == "purple":
        fill.gradient()
        fill.gradient_angle = 90.0
        fill.gradient_stops[0].color.rgb = RGBColor(156, 39, 176)
        fill.gradient_stops[1].color.rgb = RGBColor(106, 27, 154)
    elif color_theme == "blue":
        fill.gradient()
        fill.gradient_angle = 90.0
        fill.gradient_stops[0].color.rgb = RGBColor(33, 150, 243)
        fill.gradient_stops[1].color.rgb = RGBColor(21, 101, 192)
    elif color_theme == "lightgreen":
        fill.gradient()
        fill.gradient_angle = 90.0
        fill.gradient_stops[0].color.rgb = RGBColor(232, 245, 233)
        fill.gradient_stops[1].color.rgb = RGBColor(200, 230, 201)
    else:
        fill.gradient()
        fill.gradient_angle = 90.0
        fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)
        fill.gradient_stops[1].color.rgb = RGBColor(248, 249, 250)
    
    bg.line.fill.background()
    bg.element.getparent().remove(bg.element)
    slide.element.insert(0, bg.element)


def add_shadow(shape, direction=2700000):
    """Add subtle shadow to shape. Direction: 2700000=bottom, 1800000=right-bottom."""
    try:
        sp = shape.element
        spPr = sp.find('{http://schemas.openxmlformats.org/presentationml/2006/main}spPr', sp.nsmap) or sp.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        if spPr is None:
            return
        effectLst = parse_xml(f'<a:effectLst xmlns:a="{nsdecls("a")}"><a:outerShdw blurRad="50800" dist="30000" dir="{direction}" algn="ctr"><a:srgbClr val="000000"><a:alpha val="25000"/></a:srgbClr></a:outerShdw></a:effectLst>')
        spPr.append(effectLst)
    except Exception:
        pass


def add_logo(slide, prs, logo_path):
    """Add logo to top right."""
    base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    possible_paths = [
        os.path.join(base_dir, "Логотип.png"),
        logo_path,
        "Логотип.png",
    ]
    for path in possible_paths:
        if path and os.path.exists(path):
            try:
                slide.shapes.add_picture(path, prs.slide_width - Inches(2.2), Inches(0.2), height=Inches(0.75))
                return
            except Exception:
                continue


def create_donut_chart(totals, path="donut.png"):
    labels = ['Повторные\nзвонки', 'Заявки\nшт', 'Заявки\nмлн', 'Выдано\nмлн']
    values = [totals['calls_fact'], totals['leads_units_fact'], totals['leads_volume_fact']*10, totals['issued_volume']*10]
    colors = [PRIMARY, ACCENT2, '#81C784', '#AED581']
    fig = go.Figure(data=[go.Pie(
        labels=labels, 
        values=values, 
        hole=.35, 
        marker_colors=colors, 
        textposition='outside',
        textinfo='label+percent',
        textfont=dict(size=16, family="Roboto", color=TEXT_MAIN),
        pull=[0.08, 0.08, 0.08, 0.08],
        marker=dict(line=dict(color='white', width=3))
    )])
    fig.update_layout(
        title=dict(text="Распределение активности", font=dict(size=20, family="Roboto", color=TEXT_MAIN)),
        font=dict(family="Roboto", size=16, color=TEXT_MAIN), 
        showlegend=False,
        width=900, 
        height=550,
        paper_bgcolor='rgba(0,0,0,0)',
        margin=dict(l=60, r=60, t=90, b=50)
    )
    fig.write_image(path, scale=3)  # 3x DPI for ultra-sharp quality
    return path


def create_comparison_bars(prev, cur, path="compare_bars.png"):
    categories = ['Звонки', 'Заявки шт', 'Заявки млн']
    prev_vals = [prev['calls_fact'], prev['leads_units_fact'], prev['leads_volume_fact']]
    cur_vals = [cur['calls_fact'], cur['leads_units_fact'], cur['leads_volume_fact']]
    fig = go.Figure()
    fig.add_trace(go.Bar(name='Предыдущий', x=categories, y=prev_vals, marker_color=ACCENT2, text=prev_vals, textposition='outside'))
    fig.add_trace(go.Bar(name='Текущий', x=categories, y=cur_vals, marker_color=PRIMARY, text=cur_vals, textposition='outside'))
    fig.update_layout(
        barmode='group', 
        title=dict(text="Сравнение периодов", font=dict(size=18, family="Roboto")),
        font=dict(family="Roboto", size=13), 
        width=800, height=450, 
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        yaxis=dict(gridcolor='#E0E0E0')
    )
    fig.write_image(path, scale=2)
    return path


def create_line_dynamics(daily_data, path="line_dynamics.png"):
    dates = [d['date'] for d in daily_data]
    plan = [d['leads_volume_plan'] for d in daily_data]
    fact = [d['leads_volume_fact'] for d in daily_data]
    issued = [d['issued_volume'] for d in daily_data]
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=dates, y=plan, mode='lines+markers', name='План', 
                            line=dict(color=PRIMARY, width=3), marker=dict(size=8)))
    fig.add_trace(go.Scatter(x=dates, y=fact, mode='lines+markers', name='Факт', 
                            line=dict(color=ACCENT2, width=3), marker=dict(size=8)))
    fig.add_trace(go.Scatter(x=dates, y=issued, mode='lines+markers', name='Выдано', 
                            line=dict(color='#81C784', width=3), marker=dict(size=8)))
    fig.update_layout(
        title=dict(text="Динамика по дням", font=dict(size=18, family="Roboto")),
        font=dict(family="Roboto", size=13), 
        xaxis_title="Дата", yaxis_title="млн",
        width=900, height=500, 
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)',
        xaxis=dict(gridcolor='#E0E0E0'),
        yaxis=dict(gridcolor='#E0E0E0')
    )
    fig.write_image(path, scale=2)
    return path


def create_calls_line(daily_data, path="calls_line.png"):
    """Line chart for calls plan vs fact."""
    dates = [d['date'] for d in daily_data]
    plan = [d.get('calls_plan', 0) for d in daily_data]
    fact = [d.get('calls_fact', 0) for d in daily_data]
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=dates, y=plan, mode='lines+markers', name='План',
                            line=dict(color=PRIMARY, width=3), marker=dict(size=8)))
    fig.add_trace(go.Scatter(x=dates, y=fact, mode='lines+markers', name='Факт',
                            line=dict(color='#2196F3', width=3), marker=dict(size=8)))
    fig.update_layout(
        title=dict(text="Звонки: план vs факт", font=dict(size=18, family="Roboto")),
        font=dict(family="Roboto", size=13),
        xaxis_title="Дни", yaxis_title="Количество звонков",
        width=600, height=400,
        paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)',
        xaxis=dict(gridcolor='#E0E0E0'), yaxis=dict(gridcolor='#E0E0E0')
    )
    fig.write_image(path, scale=2)
    return path


def create_spider_chart(manager_data, avg_data, manager_name, path="spider.png"):
    """Radar chart: manager vs average."""
    categories = ['Повторные\nзвонки', 'Новые\nзвонки', 'Заявки\nшт', 'Заявки\nмлн', 'Одобрено', 'Выдано']
    manager_vals = [
        manager_data.calls_fact, manager_data.new_calls, manager_data.leads_units_fact,
        manager_data.leads_volume_fact, manager_data.approved_volume, manager_data.issued_volume
    ]
    avg_vals = [
        avg_data.get('calls_fact', 0), avg_data.get('new_calls', 0), avg_data.get('leads_units_fact', 0),
        avg_data.get('leads_volume_fact', 0), avg_data.get('approved_volume', 0), avg_data.get('issued_volume', 0)
    ]
    fig = go.Figure()
    fig.add_trace(go.Scatterpolar(
        r=avg_vals + [avg_vals[0]],
        theta=categories + [categories[0]],
        fill='toself',
        name='Среднее по отделу',
        fillcolor='rgba(255, 138, 101, 0.3)',
        line=dict(color=ACCENT2, width=2)
    ))
    fig.add_trace(go.Scatterpolar(
        r=manager_vals + [manager_vals[0]],
        theta=categories + [categories[0]],
        fill='toself',
        name=manager_name,
        fillcolor='rgba(156, 39, 176, 0.4)',
        line=dict(color='#9C27B0', width=3)
    ))
    fig.update_layout(
        polar=dict(radialaxis=dict(visible=True, range=[0, max(max(manager_vals), max(avg_vals)) * 1.1])),
        title=dict(text=f"Сравнение — {manager_name}", font=dict(size=16, family="Roboto")),
        font=dict(family="Roboto", size=11),
        width=550, height=450,
        paper_bgcolor='rgba(0,0,0,0)'
    )
    fig.write_image(path, scale=2)
    return path


def create_managers_bar(managers_data, path="managers_bar.png"):
    """Bar chart comparing all managers."""
    names = [m.name for m in managers_data]
    calls = [m.calls_fact for m in managers_data]
    leads = [m.leads_units_fact for m in managers_data]
    fig = go.Figure()
    fig.add_trace(go.Bar(name='Звонки', x=names, y=calls, marker_color=PRIMARY))
    fig.add_trace(go.Bar(name='Заявки', x=names, y=leads, marker_color='#2196F3'))
    fig.update_layout(
        barmode='group',
        title=dict(text="Результаты по менеджерам", font=dict(size=18, family="Roboto")),
        font=dict(family="Roboto", size=13),
        xaxis_title="Менеджеры", yaxis_title="Количество",
        width=900, height=500,
        paper_bgcolor='rgba(0,0,0,0)', plot_bgcolor='rgba(0,0,0,0)',
        yaxis=dict(gridcolor='#E0E0E0')
    )
    fig.write_image(path, scale=2)
    return path


class PremiumPresentationService:
    """Service for generating premium 9-slide PPTX presentations with charts."""
    
    def __init__(self, settings: Settings):
        self.settings = settings
        self.gpt_service = YandexGPTService(settings)
    
    async def generate_presentation(
        self,
        period_data: Dict[str, ManagerData],
        period_name: str,
        start_date: date,
        end_date: date,
        previous_data: Optional[Dict[str, ManagerData]] = None,
        previous_start_date: Optional[date] = None,
        previous_end_date: Optional[date] = None,
        daily_series: Optional[List[Dict[str, float]]] = None,
    ) -> bytes:
        """Generate premium 9-slide PPTX with charts, diagrams, AI analysis."""
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        margin = Inches(1)
        logo = self.settings.pptx_logo_path
        
        # Calculate totals
        totals = self._calculate_totals(period_data)
        prev_totals = self._calculate_totals(previous_data) if previous_data else {}
        avg = self._calculate_average_manager(period_data)
        
        # 1. Title
        await self._add_title_slide(prs, period_name, start_date, end_date, logo, margin)
        
        # 2. Team summary with table
        await self._add_team_summary_slide(prs, totals, avg, period_name, logo, margin)
        
        # 3. AI comment slide
        await self._add_ai_comment_slide(prs, totals, period_name, logo, margin)
        
        # 4. Comparison with charts
        await self._add_comparison_slide(prs, prev_totals, totals, daily_series or [], logo, margin)
        
        # 5. TOP/AntiTOP ranking
        await self._add_ranking_slide(prs, period_data, logo, margin)
        
        # 6. All managers table
        await self._add_all_managers_table(prs, period_data, logo, margin)
        
        # 7. Manager cards (2x2 grid)
        await self._add_manager_cards(prs, period_data, logo, margin)
        
        # 8. Calls dynamics (weekly line chart) - GREEN theme
        await self._add_calls_dynamics_slide(prs, daily_series or [], totals, logo, margin)
        
        # 10. Spider/Radar chart - PURPLE theme
        if period_data:
            top_manager = max(period_data.values(), key=lambda m: (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0)
            await self._add_spider_slide(prs, top_manager, avg, logo, margin)
        
        # 11. Bar chart comparison - BLUE theme
        await self._add_managers_bar_slide(prs, period_data, logo, margin)
        
        # 12. Conclusions
        await self._add_conclusions_slide(prs, totals, period_name, logo, margin)
        
        # Save
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.getvalue()
    
    async def _add_title_slide(self, prs, period_name, start_date, end_date, logo, margin):
        """Slide 1: Premium title with decorative elements."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        # Decorative top accent bar
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
        accent_bar.line.fill.background()
        
        # Large decorative frame around center
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.8), Inches(9.33), Inches(4))
        frame.fill.background()
        frame.line.color.rgb = hex_to_rgb(PRIMARY)
        frame.line.width = Pt(3)
        add_shadow(frame)
        
        # Office name - smaller, above main title
        office_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.2), Inches(8.33), Inches(0.6))
        office_box.text_frame.text = self.settings.office_name.upper()
        office_box.text_frame.paragraphs[0].font.name = "Roboto"
        office_box.text_frame.paragraphs[0].font.size = Pt(20)
        office_box.text_frame.paragraphs[0].font.bold = True
        office_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        office_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Main title
        title = slide.shapes.add_textbox(Inches(2.5), Inches(3), Inches(8.33), Inches(1.2))
        title.text_frame.text = "Отчет по продажам"
        title.text_frame.paragraphs[0].font.name = "Roboto"
        title.text_frame.paragraphs[0].font.size = Pt(54)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(TEXT_MAIN)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Divider line
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.3), Inches(4.5), Pt(2))
        divider.fill.solid()
        divider.fill.fore_color.rgb = hex_to_rgb(ACCENT2)
        divider.line.fill.background()
        
        # Period subtitle
        subtitle = slide.shapes.add_textbox(Inches(2.5), Inches(4.6), Inches(8.33), Inches(0.9))
        subtitle.text_frame.text = f"{period_name}\n{start_date.strftime('%d.%m.%Y')} — {end_date.strftime('%d.%m.%Y')}"
        for p in subtitle.text_frame.paragraphs:
            p.font.name = "Roboto"
            p.font.size = Pt(22)
            p.font.color.rgb = hex_to_rgb(TEXT_MUTED)
            p.alignment = PP_ALIGN.CENTER
    
    async def _add_team_summary_slide(self, prs, totals, avg, period_name, logo, margin):
        """Slide 2: Team summary table with zebra and traffic light."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        h = slide.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "Общие показатели команды"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(28)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Table — more compact
        rows, cols = 7, 4
        tbl = slide.shapes.add_table(rows, cols, margin, Inches(1.3), Inches(11.33), Inches(2.7)).table
        headers = ["Показатель", "План", "Факт", "Конв (%)"]
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
        
        data_rows = [
            ("Повторные звонки", f"{int(totals['calls_plan']):,}".replace(",", " "), f"{int(totals['calls_fact']):,}".replace(",", " "), f"{totals['calls_percentage']:.1f}%".replace(".", ",")),
            ("Заявки, шт", f"{int(totals['leads_units_plan']):,}".replace(",", " "), f"{int(totals['leads_units_fact']):,}".replace(",", " "), f"{totals['leads_units_percentage']:.1f}%".replace(".", ",")),
            ("Заявки, млн", f"{totals['leads_volume_plan']:.1f}".replace(".", ","), f"{totals['leads_volume_fact']:.1f}".replace(".", ","), f"{totals['leads_volume_percentage']:.1f}%".replace(".", ",")),
            ("Одобрено, млн", "—", f"{totals['approved_volume']:.1f}".replace(".", ","), "—"),
            ("Выдано, млн", "—", f"{totals['issued_volume']:.1f}".replace(".", ","), "—"),
            ("Новые звонки", "—", f"{int(totals['new_calls']):,}".replace(",", " "), "—"),
        ]
        
        for r, (name, plan, fact, conv) in enumerate(data_rows, start=1):
            for c, val in enumerate([name, plan, fact, conv]):
                cell = tbl.cell(r, c)
                cell.text = val
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Roboto"
                    p.font.size = Pt(12)
                    p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                    # Traffic light for conversion column
                    if c == 3 and conv not in ("—", "-"):
                        try:
                            pct = float(conv.replace("%", "").replace(",", "."))
                            if pct >= 90:
                                p.font.color.rgb = hex_to_rgb(PRIMARY)
                            elif pct >= 70:
                                p.font.color.rgb = hex_to_rgb(ACCENT2)
                            else:
                                p.font.color.rgb = hex_to_rgb(ALERT)
                        except Exception:
                            pass
        
        # Avg manager — right below table
        if avg:
            avg_box = slide.shapes.add_textbox(margin, Inches(4.1), Inches(11.33), Inches(0.35))
            avg_box.text_frame.text = f"📊 Средний менеджер: звонки {avg.get('calls_percentage', 0):.0f}%, заявки {avg.get('leads_volume_percentage', 0):.0f}%"
            for p in avg_box.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(11)
                p.font.italic = True
                p.font.color.rgb = hex_to_rgb(TEXT_MUTED)
        
        # Donut — full width, crisp and large
        donut_path = create_donut_chart(totals, "donut_metrics.png")
        if os.path.exists(donut_path):
            slide.shapes.add_picture(donut_path, Inches(2.2), Inches(4.7), width=Inches(9), height=Inches(2.7))
    
    async def _add_ai_comment_slide(self, prs, totals, period_name, logo, margin):
        """Slide 3: AI analysis with premium card."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        # Header with icon
        h = slide.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "🤖 Анализ и рекомендации"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Premium card for AI text — maximum size
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(6))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
        card.line.color.rgb = hex_to_rgb(PRIMARY)
        card.line.width = Pt(2)
        add_shadow(card)
        
        # AI comment inside card — comfortable padding
        ai_comment = await self.gpt_service.generate_team_comment(totals, period_name)
        ai_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.9), Inches(5.4))
        ai_box.text_frame.text = ai_comment
        ai_box.text_frame.word_wrap = True
        ai_box.text_frame.margin_left = Pt(12)
        ai_box.text_frame.margin_right = Pt(12)
        ai_box.text_frame.margin_top = Pt(12)
        ai_box.text_frame.margin_bottom = Pt(12)
        for p in ai_box.text_frame.paragraphs:
            p.font.name = "Roboto"
            p.font.size = Pt(14)
            p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
            p.line_spacing = 1.25
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)
    
    async def _add_comparison_slide(self, prs, prev_totals, cur_totals, daily_data, logo, margin):
        """Slide 4: Comparison with charts."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        h = slide.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "Сравнение с предыдущим периодом"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(28)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if prev_totals:
            compare_path = create_comparison_bars(prev_totals, cur_totals, "comparison_bars.png")
            if os.path.exists(compare_path):
                slide.shapes.add_picture(compare_path, Inches(1), Inches(1.5), width=Inches(5.5), height=Inches(3))
        
        if daily_data:
            line_path = create_line_dynamics(daily_data, "dynamics_line.png")
            if os.path.exists(line_path):
                slide.shapes.add_picture(line_path, Inches(7), Inches(1.5), width=Inches(5.5), height=Inches(3))
    
    async def _add_ranking_slide(self, prs, period_data, logo, margin):
        """Slide 5: Ranking table (simple, no overlapping shapes)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        h = slide.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "🏆 Рейтинг эффективности"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        scored = []
        for m in period_data.values():
            cp = (m.calls_fact/m.calls_plan*100) if m.calls_plan else 0
            vp = (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0
            score = 0.5*cp+0.5*vp
            scored.append((score, m.name, cp, vp))
        scored.sort(reverse=True)
        
        # Table: Rank, Name, Calls%, Volume%, Score
        rows = min(len(scored) + 1, 7)
        cols = 5
        tbl = slide.shapes.add_table(rows, cols, margin, Inches(1.5), Inches(11.33), Inches(5.5)).table
        
        headers = ["#", "Менеджер", "Звонки %", "Заявки %", "Счёт"]
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
        
        for r, (score, name, cp, vp) in enumerate(scored[:rows-1], start=1):
            row_data = [str(r), name, f"{cp:.0f}%", f"{vp:.0f}%", f"{score:.0f}"]
            for c, val in enumerate(row_data):
                cell = tbl.cell(r, c)
                cell.text = val
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
                # Traffic light for top-3 and bottom-3
                if r <= 3:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(200, 230, 201)  # light green
                elif r >= rows - 3:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 205, 210)  # light red
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Roboto"
                    p.font.size = Pt(12)
                    p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                    p.alignment = PP_ALIGN.CENTER if c > 1 else (PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)
    
    async def _add_all_managers_table(self, prs, period_data, logo, margin):
        """Slide 6: Table of all managers."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        h = slide.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "Общая таблица по менеджерам"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(28)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        rows = min(len(period_data) + 1, 8)
        cols = 5
        tbl = slide.shapes.add_table(rows, cols, margin, Inches(1.5), Inches(11.33), Inches(5)).table
        headers = ["Менеджер", "План млн", "Факт млн", "Выдано млн", "Конв %"]
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
        
        for r, (name, m) in enumerate(list(period_data.items())[:rows-1], start=1):
            conv_pct = (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0
            row_data = [name, f"{m.leads_volume_plan:.1f}".replace(".", ","), 
                       f"{m.leads_volume_fact:.1f}".replace(".", ","),
                       f"{m.issued_volume:.1f}".replace(".", ","),
                       f"{conv_pct:.1f}%".replace(".", ",")]
            for c, val in enumerate(row_data):
                cell = tbl.cell(r, c)
                cell.text = val
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Roboto"
                    p.font.size = Pt(11)
                    p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    
    async def _add_manager_cards(self, prs, period_data, logo, margin):
        """Slide 7: Manager table (simple table to avoid shape bleed)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        h = slide.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "👤 Показатели менеджеров"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Simple table instead of cards
        rows = min(len(period_data) + 1, 7)
        cols = 5
        tbl = slide.shapes.add_table(rows, cols, margin, Inches(1.5), Inches(11.33), Inches(5.5)).table
        
        headers = ["Менеджер", "План млн", "Факт млн", "Выдано млн", "Выполнение %"]
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
        
        for r, (name, m) in enumerate(list(period_data.items())[:rows-1], start=1):
            vol_pct = (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0
            row_data = [name, f"{m.leads_volume_plan:.1f}", f"{m.leads_volume_fact:.1f}", 
                       f"{m.issued_volume:.1f}", f"{vol_pct:.1f}%"]
            for c, val in enumerate(row_data):
                cell = tbl.cell(r, c)
                cell.text = val
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Roboto"
                    p.font.size = Pt(12)
                    p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    
    async def _add_calls_dynamics_slide(self, prs, daily_data, totals, logo, margin):
        """Slide 9: Calls dynamics with summary card - GREEN theme."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs, color_theme="lightgreen")
        add_logo(slide, prs, logo)
        
        # Green header bar (full width)
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = hex_to_rgb(PRIMARY)
        header_bar.line.fill.background()
        
        h = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "📞 ДИНАМИКА ЗВОНКОВ (НЕДЕЛЯ)"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Line chart
        if daily_data:
            calls_path = create_calls_line(daily_data, "calls_weekly.png")
            if os.path.exists(calls_path):
                slide.shapes.add_picture(calls_path, Inches(1.5), Inches(1.8), width=Inches(6.5), height=Inches(4.5))
        
        # Summary card
        summary_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2), Inches(4), Inches(4.5))
        summary_card.fill.solid()
        summary_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        summary_card.line.color.rgb = hex_to_rgb(PRIMARY)
        summary_card.line.width = Pt(2)
        add_shadow(summary_card)
        
        summary_text = (
            f"📊 ИТОГИ НЕДЕЛИ\n\n"
            f"🎯 План звонков: {int(totals.get('calls_plan', 0)):,}\n"
            f"✅ Выполнено: {int(totals.get('calls_fact', 0)):,}\n"
            f"📊 Процент: {totals.get('calls_percentage', 0):.1f}%\n\n"
            f"💡 Новые контакты: {int(totals.get('new_calls', 0)):,}"
        )
        summary_box = slide.shapes.add_textbox(Inches(8.8), Inches(2.3), Inches(3.4), Inches(4))
        summary_box.text_frame.text = summary_text
        for p in summary_box.text_frame.paragraphs:
            p.font.name = "Roboto"
            p.font.size = Pt(13)
            p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
            p.space_after = Pt(6)
    
    async def _add_spider_slide(self, prs, manager, avg, logo, margin):
        """Slide 10: Spider/Radar chart - PURPLE theme."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs, color_theme="purple")
        
        # Purple header bar
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = RGBColor(156, 39, 176)
        header_bar.line.fill.background()
        
        h = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = f"📡 ПРОФИЛЬ ЭФФЕКТИВНОСТИ"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle - show manager name only once
        sub = slide.shapes.add_textbox(margin, Inches(1.5), prs.slide_width - 2*margin, Inches(0.4))
        sub.text_frame.text = f"{manager.name} vs Средний менеджер"
        sub.text_frame.paragraphs[0].font.name = "Roboto"
        sub.text_frame.paragraphs[0].font.size = Pt(20)
        sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Spider chart
        spider_path = create_spider_chart(manager, avg, manager.name, "spider_chart.png")
        if os.path.exists(spider_path):
            slide.shapes.add_picture(spider_path, Inches(3.5), Inches(2.2), width=Inches(6.5), height=Inches(5))
    
    async def _add_managers_bar_slide(self, prs, period_data, logo, margin):
        """Slide 11: Bar chart - BLUE theme."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs, color_theme="blue")
        
        # Blue header bar
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = RGBColor(33, 150, 243)
        header_bar.line.fill.background()
        
        h = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "📊 СРАВНЕНИЕ КОМАНДЫ"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Bar chart
        managers_list = list(period_data.values())
        if managers_list:
            bar_path = create_managers_bar(managers_list, "managers_comparison.png")
            if os.path.exists(bar_path):
                slide.shapes.add_picture(bar_path, Inches(2), Inches(1.8), width=Inches(9.33), height=Inches(5))
    
    async def _add_conclusions_slide(self, prs, totals, period_name, logo, margin):
        """Slide 9: AI conclusions with premium card."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_gradient_bg(slide, prs)
        add_logo(slide, prs, logo)
        
        # Header
        h = slide.shapes.add_textbox(margin, Inches(0.5), prs.slide_width - 2*margin, Inches(0.6))
        h.text_frame.text = "✅ Выводы и рекомендации"
        h.text_frame.paragraphs[0].font.name = "Roboto"
        h.text_frame.paragraphs[0].font.size = Pt(32)
        h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        h.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Premium card — maximum size, shadow to right-bottom to avoid overflow
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.9))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(CARD_BG)
        card.line.color.rgb = hex_to_rgb(PRIMARY)
        card.line.width = Pt(2)
        add_shadow(card, direction=1800000)
        
        ai_conclusion = await self.gpt_service.generate_team_comment(totals, f"Итоги: {period_name}")
        ai_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.9), Inches(5.4))
        ai_box.text_frame.text = ai_conclusion
        ai_box.text_frame.word_wrap = True
        ai_box.text_frame.margin_left = Pt(12)
        ai_box.text_frame.margin_right = Pt(12)
        ai_box.text_frame.margin_top = Pt(12)
        ai_box.text_frame.margin_bottom = Pt(12)
        for p in ai_box.text_frame.paragraphs:
            p.font.name = "Roboto"
            p.font.size = Pt(14)
            p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
            p.line_spacing = 1.25
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)
    
    def _calculate_totals(self, period_data: Dict[str, ManagerData]) -> Dict[str, float]:
        """Calculate team totals."""
        if not period_data:
            return {}
        totals = {
            'calls_plan': 0, 'calls_fact': 0,
            'leads_units_plan': 0, 'leads_units_fact': 0,
            'leads_volume_plan': 0.0, 'leads_volume_fact': 0.0,
            'approved_volume': 0.0, 'issued_volume': 0.0,
            'new_calls': 0, 'new_calls_plan': 0,
        }
        for m in period_data.values():
            totals['calls_plan'] += m.calls_plan
            totals['calls_fact'] += m.calls_fact
            totals['leads_units_plan'] += m.leads_units_plan
            totals['leads_units_fact'] += m.leads_units_fact
            totals['leads_volume_plan'] += m.leads_volume_plan
            totals['leads_volume_fact'] += m.leads_volume_fact
            totals['approved_volume'] += m.approved_volume
            totals['issued_volume'] += m.issued_volume
            totals['new_calls'] += m.new_calls
            totals['new_calls_plan'] += m.new_calls_plan
        totals['calls_percentage'] = (totals['calls_fact'] / totals['calls_plan'] * 100) if totals['calls_plan'] else 0
        totals['leads_units_percentage'] = (totals['leads_units_fact'] / totals['leads_units_plan'] * 100) if totals['leads_units_plan'] else 0
        totals['leads_volume_percentage'] = (totals['leads_volume_fact'] / totals['leads_volume_plan'] * 100) if totals['leads_volume_plan'] else 0
        return totals
    
    def _calculate_average_manager(self, period_data: Dict[str, ManagerData]) -> Dict[str, float]:
        """Calculate average manager baseline."""
        if not period_data:
            return {}
        n = len(period_data)
        avg = {'calls_plan': 0, 'calls_fact': 0, 'leads_units_plan': 0, 'leads_units_fact': 0,
               'leads_volume_plan': 0.0, 'leads_volume_fact': 0.0,
               'approved_volume': 0.0, 'issued_volume': 0.0, 'new_calls': 0}
        for m in period_data.values():
            avg['calls_plan'] += m.calls_plan
            avg['calls_fact'] += m.calls_fact
            avg['leads_units_plan'] += m.leads_units_plan
            avg['leads_units_fact'] += m.leads_units_fact
            avg['leads_volume_plan'] += m.leads_volume_plan
            avg['leads_volume_fact'] += m.leads_volume_fact
            avg['approved_volume'] += m.approved_volume
            avg['issued_volume'] += m.issued_volume
            avg['new_calls'] += m.new_calls
        for k in avg:
            avg[k] = avg[k] / n
        avg['calls_percentage'] = (avg['calls_fact'] / avg['calls_plan'] * 100) if avg['calls_plan'] else 0
        avg['leads_units_percentage'] = (avg['leads_units_fact'] / avg['leads_units_plan'] * 100) if avg['leads_units_plan'] else 0
        avg['leads_volume_percentage'] = (avg['leads_volume_fact'] / avg['leads_volume_plan'] * 100) if avg['leads_volume_plan'] else 0
        return avg

