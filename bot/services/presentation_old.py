"""Presentation generation service."""
import os
import io
from datetime import datetime, date
from typing import Dict, Any, List, Tuple, Optional
from dataclasses import dataclass

import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Cm

from bot.services.yandex_gpt import YandexGPTService
from bot.config import Settings


@dataclass
class ManagerData:
    """Data structure for manager statistics."""
    name: str
    calls_plan: int = 0
    calls_fact: int = 0
    leads_units_plan: int = 0
    leads_units_fact: int = 0
    leads_volume_plan: float = 0.0
    leads_volume_fact: float = 0.0
    approved_volume: float = 0.0
    issued_volume: float = 0.0
    new_calls: int = 0
    new_calls_plan: int = 0
    
    @property
    def calls_percentage(self) -> float:
        """Calculate calls completion percentage."""
        return (self.calls_fact / self.calls_plan * 100) if self.calls_plan > 0 else 0
    
    @property
    def leads_units_percentage(self) -> float:
        """Calculate leads units completion percentage."""
        return (self.leads_units_fact / self.leads_units_plan * 100) if self.leads_units_plan > 0 else 0
    
    @property
    def leads_volume_percentage(self) -> float:
        """Calculate leads volume completion percentage."""
        return (self.leads_volume_fact / self.leads_volume_plan * 100) if self.leads_volume_plan > 0 else 0


class PresentationService:
    """Service for generating PowerPoint presentations."""
    
    def __init__(self, settings: Settings):
        self.settings = settings
        self.gpt_service = YandexGPTService(settings)
    
    # Helpers: branding and colors
    def _rgb_from_hex(self, hex_color: str) -> RGBColor:
        try:
            r = int(hex_color[1:3], 16)
            g = int(hex_color[3:5], 16)
            b = int(hex_color[5:7], 16)
            return RGBColor(r, g, b)
        except Exception:
            return RGBColor(204, 0, 0)
    
    def _apply_brand(self, slide) -> None:
        try:
            pres = slide.part.presentation
            width = pres.slide_width
            # Top band background
            band = slide.shapes.add_shape(
                1,  # MSO_AUTO_SHAPE_TYPE = Rectangle
                0, 0, width, Inches(0.6)
            )
            band.fill.solid()
            band.fill.fore_color.rgb = self._rgb_from_hex(getattr(self.settings, 'pptx_secondary_color', '#F3F4F6'))
            band.line.fill.background()
            # Logo (optional)
            if getattr(self.settings, 'pptx_logo_path', '') and os.path.exists(self.settings.pptx_logo_path):
                slide.shapes.add_picture(
                    self.settings.pptx_logo_path,
                    width - Inches(1.8), Inches(0.1), height=Inches(0.4)
                )
        except Exception:
            pass
    
    async def generate_presentation(
        self,
        period_data: Dict[str, ManagerData],
        period_name: str,
        start_date: date,
        end_date: date,
        previous_data: Optional[Dict[str, ManagerData]] = None,
        previous_start_date: Optional[date] = None,
        previous_end_date: Optional[date] = None,
    ) -> bytes:
        """
        Generate PowerPoint presentation with analytics.
        
        Args:
            period_data: Dictionary mapping manager names to their data
            period_name: Human-readable period name (e.g., "Неделя 18-24 августа")
            start_date: Period start date
            end_date: Period end date
            
        Returns:
            PPTX file as bytes
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide size (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Apply optional logo on master (top-right) and light gray background band
        try:
            if self.settings.pptx_logo_path and os.path.exists(self.settings.pptx_logo_path):
                for layout in prs.slide_layouts:
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.add_picture(self.settings.pptx_logo_path, prs.slide_width - Inches(1.8), Inches(0.2), height=Inches(0.9))
                    # remove after cloning: keep normal slides clean
                    prs.slides._sldIdLst.remove(slide._element.getparent())
        except Exception:
            pass
        
        # Title slide
        await self._add_title_slide(prs, period_name, start_date, end_date)
        
        # Summary slide
        await self._add_summary_slide(prs, period_data, period_name)

        # Comparison slide (previous vs current)
        if previous_data is not None:
            await self._add_comparison_slide(
                prs,
                previous_data,
                period_data,
                start_date,
                end_date,
                previous_start_date,
                previous_end_date,
            )
        
        # Per‑manager: only comparison slide (tables + AI‑комментарий), без отдельной страницы с показателями
        for manager_name, manager_data in period_data.items():
            if previous_data is not None and manager_name in previous_data:
                await self._add_manager_comparison_slide(
                    prs,
                    previous_data[manager_name],
                    manager_data,
                    start_date,
                    end_date,
                    previous_start_date,
                    previous_end_date,
                    period_name,
                )
        
        # Team AI analysis slide is omitted per revised presentation flow
        
        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        return pptx_buffer.getvalue()
    
    async def _add_title_slide(
        self,
        prs: Presentation,
        period_name: str,
        start_date: date,
        end_date: date
    ):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        self._apply_brand(slide)
        
        # Title
        title = slide.shapes.title
        title.text = f"Отчет по продажам"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        # Primary color
        try:
            r = int(self.settings.pptx_primary_color[1:3], 16)
            g = int(self.settings.pptx_primary_color[3:5], 16)
            b = int(self.settings.pptx_primary_color[5:7], 16)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(r, g, b)
        except Exception:
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Force full-width title box for perfect centering
        try:
            title.left = 0
            title.width = prs.slide_width
        except Exception:
            pass
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"{period_name}\n{start_date.strftime('%d.%m.%Y')} — {end_date.strftime('%d.%m.%Y')}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        subtitle.text_frame.paragraphs[1].font.size = Pt(20)
        subtitle.text_frame.paragraphs[1].font.name = self.settings.pptx_font_family
        subtitle.text_frame.paragraphs[1].font.color.rgb = RGBColor(102, 102, 102)  # Gray
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle.text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    async def _add_summary_slide(
        self,
        prs: Presentation,
        period_data: Dict[str, ManagerData],
        period_name: str
    ):
        """Add summary slide with team totals."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        self._apply_brand(slide)
        
        # Title
        title = slide.shapes.title
        title.text = f"Общие показатели команды"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        try:
            r = int(self.settings.pptx_primary_color[1:3], 16)
            g = int(self.settings.pptx_primary_color[3:5], 16)
            b = int(self.settings.pptx_primary_color[5:7], 16)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(r, g, b)
        except Exception:
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        try:
            title.left = 0
            title.width = prs.slide_width
        except Exception:
            pass
        
        # Calculate totals
        totals = self._calculate_totals(period_data)
        
        # Content placeholder: clear to avoid overlap, we will render a table instead
        content = slide.placeholders[1]
        try:
            content.text = ""
        except Exception:
            pass

        # Calculate average manager baseline
        avg = self._calculate_average_manager(period_data)

        # Summary table
        top = Inches(1.8)
        left = Inches(0.5)
        width = Inches(12.3)
        height = Inches(2.6)
        rows = 7  # header + 6 metrics
        cols = 4  # metric, plan, fact, conv
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        headers = ["Показатель", "План", "Факт", "Конв (%)"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = self.settings.pptx_font_family
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb_from_hex(self.settings.pptx_primary_color)
                p.font.color.rgb = RGBColor(255, 255, 255)
            except Exception:
                pass
            p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT

        def set_row(r: int, name: str, plan: str, fact: str, conv: str) -> None:
            values = [name, plan, fact, conv]
            for c, val in enumerate(values):
                table.cell(r, c).text = val
                pp = table.cell(r, c).text_frame.paragraphs[0]
                pp.font.size = Pt(12)
                pp.font.name = self.settings.pptx_font_family
                pp.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

        # Compute conversions
        calls_conv = f"{totals['calls_percentage']:.1f}%" if totals['calls_plan'] else "-"
        units_conv = f"{totals['leads_units_percentage']:.1f}%" if totals['leads_units_plan'] else "-"
        vol_conv = f"{totals['leads_volume_percentage']:.1f}%" if totals['leads_volume_plan'] else "-"

        # Traffic-light fill helper
        def conv_rgb(value_str: str) -> RGBColor | None:
            try:
                v = float(value_str.replace('%', ''))
            except Exception:
                return None
            if v >= 90:
                return RGBColor(46, 125, 50)  # green
            if v >= 70:
                return RGBColor(255, 138, 101)  # amber
            return RGBColor(198, 40, 40)  # red

        # Fill rows
        set_row(1, "📲 Повторные звонки", f"{totals['calls_plan']:,}", f"{totals['calls_fact']:,}", calls_conv)
        set_row(2, "📝 Заявки, шт", f"{totals['leads_units_plan']:,}", f"{totals['leads_units_fact']:,}", units_conv)
        set_row(3, "💰 Заявки, млн", f"{totals['leads_volume_plan']:.1f}", f"{totals['leads_volume_fact']:.1f}", vol_conv)
        set_row(4, "✅ Одобрено, млн", "-", f"{totals['approved_volume']:.1f}", "-")
        set_row(5, "✅ Выдано, млн", "-", f"{totals['issued_volume']:.1f}", "-")
        set_row(6, "☎️ Новые звонки", "-", f"{totals['new_calls']:,}", "-")

        # Apply traffic-light color to conversion column cells (rows 1..3, col=3)
        try:
            mapping = {1: calls_conv, 2: units_conv, 3: vol_conv}
            for r in (1, 2, 3):
                color = conv_rgb(mapping[r])
                if color:
                    cell = table.cell(r, 3)
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = color
        except Exception:
            pass

        # Add baseline comparison row below table
        # "Средний менеджер" as reference
        if avg:
            baseline_box = prs.slides[-1].shapes.add_textbox(left, top + height + Inches(0.1), width, Inches(0.5))
            bf = baseline_box.text_frame
            bf.text = f"📊 Средний менеджер: звонки {avg.get('calls_percentage', 0):.0f}%, заявки {avg.get('leads_volume_percentage', 0):.0f}%"
            bf.paragraphs[0].font.size = Pt(11)
            bf.paragraphs[0].font.name = self.settings.pptx_font_family
            bf.paragraphs[0].font.italic = True
            bf.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

        # Add AI team comment below
        # Place AI team comment right under the table (increase height to avoid clipping)
        comment_box = prs.slides[-1].shapes.add_textbox(left, top + height + Inches(0.7), width, Inches(2.5))
        tf = comment_box.text_frame
        tf.text = "Комментарий ИИ — Команда"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.name = self.settings.pptx_font_family
        tf.paragraphs[0].font.bold = True
        ai_comment = await self.gpt_service.generate_team_comment(totals, period_name)
        p = tf.add_paragraph()
        p.text = ai_comment
        p.font.size = Pt(12)
        p.font.name = self.settings.pptx_font_family
        # Tighten spacing to avoid overflow
        for par in tf.paragraphs:
            try:
                par.space_after = Pt(4)
            except Exception:
                pass
    
    async def _add_manager_slide(self, prs: Presentation, manager_data: ManagerData):
        """Add individual manager slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        self._apply_brand(slide)
        
        # Title
        title = slide.shapes.title
        title.text = f"👤 {manager_data.name}"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        try:
            r = int(self.settings.pptx_primary_color[1:3], 16)
            g = int(self.settings.pptx_primary_color[3:5], 16)
            b = int(self.settings.pptx_primary_color[5:7], 16)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(r, g, b)
        except Exception:
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        try:
            title.left = 0
            title.width = prs.slide_width
        except Exception:
            pass
        
        # Performance indicators
        calls_status = "🟢" if manager_data.calls_percentage >= 80 else "🟡" if manager_data.calls_percentage >= 60 else "🔴"
        leads_status = "🟢" if manager_data.leads_volume_percentage >= 80 else "🟡" if manager_data.leads_volume_percentage >= 60 else "🔴"
        
        # Content
        content = slide.placeholders[1]
        content.text = f"""📈 Показатели эффективности

{calls_status} Повторные звонки: {manager_data.calls_fact:,} из {manager_data.calls_plan:,} ({manager_data.calls_percentage:.1f}%)

📝 Заявки (шт): {manager_data.leads_units_fact:,} из {manager_data.leads_units_plan:,} ({manager_data.leads_units_percentage:.1f}%)

{leads_status} Заявки (млн): {manager_data.leads_volume_fact:.1f} из {manager_data.leads_volume_plan:.1f} ({manager_data.leads_volume_percentage:.1f}%)

✅ Одобрено (млн): {manager_data.approved_volume:.1f}

✅ Выдано (млн): {manager_data.issued_volume:.1f}

☎️ Новые звонки: {manager_data.new_calls:,}"""
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.name = self.settings.pptx_font_family
            paragraph.space_after = Pt(8)

    async def _add_manager_comparison_slide(
        self,
        prs: Presentation,
        prev: ManagerData,
        cur: ManagerData,
        current_start: date,
        current_end: date,
        previous_start: Optional[date],
        previous_end: Optional[date],
        period_name: str,
    ) -> None:
        """Add per-manager comparison slide with two tables + totals + AI comment on one slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        self._apply_brand(slide)
        title = slide.shapes.title
        title.text = f"Динамика — {cur.name}"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        try:
            r = int(self.settings.pptx_primary_color[1:3], 16)
            g = int(self.settings.pptx_primary_color[3:5], 16)
            b = int(self.settings.pptx_primary_color[5:7], 16)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(r, g, b)
        except Exception:
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        try:
            title.left = 0
            title.width = prs.slide_width
        except Exception:
            pass

        def t(m: ManagerData) -> dict[str, float]:
            return {
                'calls_plan': m.calls_plan,
                'calls_fact': m.calls_fact,
                'leads_units_plan': m.leads_units_plan,
                'leads_units_fact': m.leads_units_fact,
                'leads_volume_plan': m.leads_volume_plan,
                'leads_volume_fact': m.leads_volume_fact,
                'approved_volume': m.approved_volume,
                'issued_volume': m.issued_volume,
                'new_calls': m.new_calls,
            }

        prev_d = t(prev)
        cur_d = t(cur)

        rows = 5
        cols = 4
        left_prev = Inches(0.5)
        top_prev = Inches(1.6)
        width = Inches(6.0)
        height = Inches(2.3)
        table_prev = slide.shapes.add_table(rows, cols, left_prev, top_prev, width, height).table
        table_cur = slide.shapes.add_table(rows, cols, left_prev + Inches(6.3), top_prev, width, height).table

        # Add period captions above tables
        if previous_start and previous_end:
            cap_prev = slide.shapes.add_textbox(left_prev, top_prev - Inches(0.35), width, Inches(0.3))
            cp = cap_prev.text_frame
            cp.text = f"Предыдущий период: {previous_start.strftime('%d.%m.%Y')} — {previous_end.strftime('%d.%m.%Y')}"
            cp.paragraphs[0].font.size = Pt(12)
            cp.paragraphs[0].font.name = self.settings.pptx_font_family
            cp.paragraphs[0].alignment = PP_ALIGN.CENTER
        cap_cur = slide.shapes.add_textbox(left_prev + Inches(6.3), top_prev - Inches(0.35), width, Inches(0.3))
        cc = cap_cur.text_frame
        cc.text = f"Текущий период: {current_start.strftime('%d.%m.%Y')} — {current_end.strftime('%d.%m.%Y')}"
        cc.paragraphs[0].font.size = Pt(12)
        cc.paragraphs[0].font.name = self.settings.pptx_font_family
        cc.paragraphs[0].alignment = PP_ALIGN.CENTER

        headers = ["Показатель", "План", "Факт", "Конв (%)"]
        for i, h in enumerate(headers):
            for tbl in (table_prev, table_cur):
                cell = tbl.cell(0, i)
                cell.text = h
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.name = self.settings.pptx_font_family
                # ensure emojis render on platforms by setting emoji-capable font for header
                try:
                    p.font.name = self.settings.pptx_font_family
                except Exception:
                    pass
                # Center Plan/Fact/Conv headers
                p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
                # Header background tint
                try:
                    r = int(self.settings.pptx_primary_color[1:3], 16)
                    g = int(self.settings.pptx_primary_color[3:5], 16)
                    b = int(self.settings.pptx_primary_color[5:7], 16)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(r, g, b)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                except Exception:
                    pass

        metrics = [
            ("📲 Повторные звонки", 'calls_plan', 'calls_fact'),
            ("☎️ Новые звонки", 'new_calls_plan', 'new_calls'),
            ("📝 Заявки, шт", 'leads_units_plan', 'leads_units_fact'),
            ("💰 Заявки, млн", 'leads_volume_plan', 'leads_volume_fact'),
        ]

        def fill_row(tbl, row_idx, name, plan_val, fact_val):
            tbl.cell(row_idx, 0).text = name
            tbl.cell(row_idx, 1).text = f"{plan_val:,.1f}" if isinstance(plan_val, float) else f"{plan_val:,}"
            tbl.cell(row_idx, 2).text = f"{fact_val:,.1f}" if isinstance(fact_val, float) else f"{fact_val:,}"
            conv = (fact_val / plan_val * 100) if (isinstance(plan_val, (int, float)) and plan_val) else 0
            tbl.cell(row_idx, 3).text = f"{conv:.1f}%"
            for c in range(4):
                p = tbl.cell(row_idx, c).text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.name = self.settings.pptx_font_family
                # Center Plan/Fact/Conv data
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

        for idx, (name, plan_key, fact_key) in enumerate(metrics, start=1):
            fill_row(table_prev, idx, name, prev_d.get(plan_key, 0), prev_d.get(fact_key, 0))
            fill_row(table_cur, idx, name, cur_d.get(plan_key, 0), cur_d.get(fact_key, 0))

        textbox_prev = slide.shapes.add_textbox(left_prev, top_prev + Inches(2.5), width, Inches(0.9))
        tfp = textbox_prev.text_frame
        tfp.text = (
            f"💰 План: {prev_d['leads_volume_plan']:.1f} млн\n"
            f"✅ Одобрено: {prev_d['approved_volume']:.1f} млн\n"
            f"✅ Выдано: {prev_d['issued_volume']:.1f} млн\n"
            f"🎯 Осталось выдать: {max(prev_d['leads_volume_plan'] - prev_d['issued_volume'], 0):.1f} млн"
        )
        for p in tfp.paragraphs:
            p.font.size = Pt(11)
            p.font.name = self.settings.pptx_font_family

        textbox_cur = slide.shapes.add_textbox(left_prev + Inches(6.3), top_prev + Inches(2.5), width, Inches(0.9))
        tfc = textbox_cur.text_frame
        tfc.text = (
            f"💰 План: {cur_d['leads_volume_plan']:.1f} млн\n"
            f"✅ Одобрено: {cur_d['approved_volume']:.1f} млн\n"
            f"✅ Выдано: {cur_d['issued_volume']:.1f} млн\n"
            f"🎯 Осталось выдать: {max(cur_d['leads_volume_plan'] - cur_d['issued_volume'], 0):.1f} млн"
        )
        for p in tfc.paragraphs:
            p.font.size = Pt(11)
            p.font.name = self.settings.pptx_font_family

        # AI comment block on the same slide
        def as_dict(m: ManagerData) -> dict[str, float]:
            return {
                'calls_plan': m.calls_plan,
                'calls_fact': m.calls_fact,
                'leads_units_plan': m.leads_units_plan,
                'leads_units_fact': m.leads_units_fact,
                'leads_volume_plan': m.leads_volume_plan,
                'leads_volume_fact': m.leads_volume_fact,
                'approved_volume': m.approved_volume,
                'issued_volume': m.issued_volume,
                'new_calls': m.new_calls,
            }

        prev_dict = as_dict(prev)
        cur_dict = as_dict(cur)
        comment = await self.gpt_service.generate_manager_comment(cur.name, prev_dict, cur_dict, period_name)

        # Place comment higher and allow wrapping to avoid clipping on last slide
        # Place comment below totals with safe margin to avoid overlap
        comment_top = top_prev + Inches(2.5) + Inches(1.0)
        comment_box = slide.shapes.add_textbox(Inches(0.5), comment_top, Inches(12.3), Inches(3.0))
        tfc = comment_box.text_frame
        # Heading
        tfc.text = f"Комментарий ИИ — {cur.name}"
        tfc.paragraphs[0].font.size = Pt(13)
        tfc.paragraphs[0].font.name = self.settings.pptx_font_family
        tfc.paragraphs[0].font.bold = True
        # Body
        p = tfc.add_paragraph()
        p.text = comment
        p.font.size = Pt(11)
        p.font.name = self.settings.pptx_font_family
        try:
            for par in tfc.paragraphs:
                par.space_after = Pt(3)
        except Exception:
            pass

    # Backward-compatibility: older callers might still invoke this to add a separate AI comment slide
    async def _add_manager_ai_comment_slide(self, prs: Presentation, prev: ManagerData, cur: ManagerData, period_name: str) -> None:
        """Legacy method: add a standalone AI comment slide for a manager (kept for compatibility)."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        self._apply_brand(slide)
        title = slide.shapes.title
        title.text = f"Комментарий ИИ — {cur.name}"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        try:
            r = int(self.settings.pptx_primary_color[1:3], 16)
            g = int(self.settings.pptx_primary_color[3:5], 16)
            b = int(self.settings.pptx_primary_color[5:7], 16)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(r, g, b)
        except Exception:
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)

        def as_dict(m: ManagerData) -> dict[str, float]:
            return {
                'calls_plan': m.calls_plan,
                'calls_fact': m.calls_fact,
                'leads_units_plan': m.leads_units_plan,
                'leads_units_fact': m.leads_units_fact,
                'leads_volume_plan': m.leads_volume_plan,
                'leads_volume_fact': m.leads_volume_fact,
                'approved_volume': m.approved_volume,
                'issued_volume': m.issued_volume,
                'new_calls': m.new_calls,
            }

        prev_dict = as_dict(prev)
        cur_dict = as_dict(cur)
        comment = await self.gpt_service.generate_manager_comment(cur.name, prev_dict, cur_dict, period_name)

        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.5), Inches(4.5))
        tf = textbox.text_frame
        tf.text = comment
        for p in tf.paragraphs:
            p.font.size = Pt(16)
            p.font.name = self.settings.pptx_font_family
    
    async def _add_ai_analysis_slide(
        self,
        prs: Presentation,
        period_data: Dict[str, ManagerData],
        period_name: str
    ):
        """Add AI analysis slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "🤖 AI-Анализ и рекомендации"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        
        # Generate AI analysis
        analysis_data = {}
        for name, data in period_data.items():
            analysis_data[name] = {
                'calls_plan': data.calls_plan,
                'calls_fact': data.calls_fact,
                'leads_units_plan': data.leads_units_plan,
                'leads_units_fact': data.leads_units_fact,
                'leads_volume_plan': data.leads_volume_plan,
                'leads_volume_fact': data.leads_volume_fact,
                'approved_volume': data.approved_volume,
                'issued_volume': data.issued_volume,
            }
        
        ai_analysis = await self.gpt_service.generate_analysis(analysis_data)
        
        # Content
        content = slide.placeholders[1]
        content.text = ai_analysis
        
        # Format content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.name = self.settings.pptx_font_family
            paragraph.space_after = Pt(6)

    async def _add_comparison_slide(
        self,
        prs: Presentation,
        previous_data: Dict[str, ManagerData],
        current_data: Dict[str, ManagerData],
        current_start: date,
        current_end: date,
        previous_start: Optional[date],
        previous_end: Optional[date],
    ) -> None:
        """Add team comparison slide with centered header and period captions over tables."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        title = slide.shapes.title
        title.text = "Динамика: предыдущий период vs текущий"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        try:
            title.left = 0
            title.width = prs.slide_width
        except Exception:
            pass

        def totals(data: Dict[str, ManagerData]) -> Dict[str, float]:
            t = {
                'calls_plan': 0, 'calls_fact': 0,
                'leads_units_plan': 0, 'leads_units_fact': 0,
                'leads_volume_plan': 0.0, 'leads_volume_fact': 0.0,
                'approved_volume': 0.0, 'issued_volume': 0.0,
                'new_calls': 0,
                'new_calls_plan': 0,
            }
            for d in data.values():
                t['calls_plan'] += d.calls_plan
                t['calls_fact'] += d.calls_fact
                t['leads_units_plan'] += d.leads_units_plan
                t['leads_units_fact'] += d.leads_units_fact
                t['leads_volume_plan'] += d.leads_volume_plan
                t['leads_volume_fact'] += d.leads_volume_fact
                t['approved_volume'] += d.approved_volume
                t['issued_volume'] += d.issued_volume
                t['new_calls'] += d.new_calls
                t['new_calls_plan'] += d.new_calls_plan
            return t

        prev = totals(previous_data)
        cur = totals(current_data)

        # Create two tables
        rows = 5  # headers + 4 metrics
        cols = 4  # metric name + Plan + Fact + Conv
        left_prev = Inches(0.5)
        top_prev = Inches(1.8)
        width = Inches(6.0)
        height = Inches(2.5)
        table_prev = slide.shapes.add_table(rows, cols, left_prev, top_prev, width, height).table
        table_cur = slide.shapes.add_table(rows, cols, left_prev + Inches(6.3), top_prev, width, height).table

        # Period captions above tables
        if previous_start and previous_end:
            cap_prev = slide.shapes.add_textbox(left_prev, top_prev - Inches(0.35), width, Inches(0.3))
            cp = cap_prev.text_frame
            cp.text = f"Предыдущий период: {previous_start.strftime('%d.%m.%Y')} — {previous_end.strftime('%d.%m.%Y')}"
            cp.paragraphs[0].font.size = Pt(12)
            cp.paragraphs[0].font.name = self.settings.pptx_font_family
            cp.paragraphs[0].alignment = PP_ALIGN.CENTER
        cap_cur = slide.shapes.add_textbox(left_prev + Inches(6.3), top_prev - Inches(0.35), width, Inches(0.3))
        cc = cap_cur.text_frame
        cc.text = f"Текущий период: {current_start.strftime('%d.%m.%Y')} — {current_end.strftime('%d.%m.%Y')}"
        cc.paragraphs[0].font.size = Pt(12)
        cc.paragraphs[0].font.name = self.settings.pptx_font_family
        cc.paragraphs[0].alignment = PP_ALIGN.CENTER

        headers = ["Показатель", "План", "Факт", "Конв (%)"]
        for i, h in enumerate(headers):
            for tbl in (table_prev, table_cur):
                cell = tbl.cell(0, i)
                cell.text = h
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.name = self.settings.pptx_font_family
                # Center Plan/Fact/Conv headers
                p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._rgb_from_hex(self.settings.pptx_primary_color)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                except Exception:
                    pass

        metrics = [
            ("📲 Повторные звонки", 'calls_plan', 'calls_fact'),
            ("☎️ Новые звонки", 'new_calls_plan', 'new_calls'),
            ("📝 Заявки, шт", 'leads_units_plan', 'leads_units_fact'),
            ("💰 Заявки, млн", 'leads_volume_plan', 'leads_volume_fact'),
        ]

        def fill_row(tbl, row_idx, name, plan_val, fact_val):
            tbl.cell(row_idx, 0).text = name
            tbl.cell(row_idx, 1).text = f"{plan_val:,.1f}" if isinstance(plan_val, float) else f"{plan_val:,}"
            tbl.cell(row_idx, 2).text = f"{fact_val:,.1f}" if isinstance(fact_val, float) else f"{fact_val:,}"
            conv = (fact_val / plan_val * 100) if (isinstance(plan_val, (int, float)) and plan_val) else 0
            tbl.cell(row_idx, 3).text = f"{conv:.1f}%"
            for c in range(4):
                p = tbl.cell(row_idx, c).text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.name = self.settings.pptx_font_family
                # Center Plan/Fact/Conv data
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

        for idx, (name, plan_key, fact_key) in enumerate(metrics, start=1):
            fill_row(table_prev, idx, name, prev.get(plan_key, 0), prev.get(fact_key, 0))
            fill_row(table_cur, idx, name, cur.get(plan_key, 0), cur.get(fact_key, 0))

        # Totals summary text boxes below tables
        textbox_prev = slide.shapes.add_textbox(left_prev, top_prev + Inches(2.7), width, Inches(1.2))
        tfp = textbox_prev.text_frame
        tfp.text = (
            f"💰 План: {prev['leads_volume_plan']:.1f} млн\n"
            f"✅ Одобрено: {prev['approved_volume']:.1f} млн\n"
            f"✅ Выдано: {prev['issued_volume']:.1f} млн\n"
            f"🎯 Осталось выдать: {max(prev['leads_volume_plan'] - prev['issued_volume'], 0):.1f} млн"
        )
        for p in tfp.paragraphs:
            p.font.size = Pt(12)
            p.font.name = self.settings.pptx_font_family

        textbox_cur = slide.shapes.add_textbox(left_prev + Inches(6.3), top_prev + Inches(2.7), width, Inches(1.2))
        tfc = textbox_cur.text_frame
        tfc.text = (
            f"💰 План: {cur['leads_volume_plan']:.1f} млн\n"
            f"✅ Одобрено: {cur['approved_volume']:.1f} млн\n"
            f"✅ Выдано: {cur['issued_volume']:.1f} млн\n"
            f"🎯 Осталось выдать: {max(cur['leads_volume_plan'] - cur['issued_volume'], 0):.1f} млн"
        )
        for p in tfc.paragraphs:
            p.font.size = Pt(12)
            p.font.name = self.settings.pptx_font_family

        # AI comparison comment block
        def totals_dict(data: Dict[str, float]) -> Dict[str, float]:
            return data

        comment_top = top_prev + Inches(3.9)
        comment_box = slide.shapes.add_textbox(Inches(0.5), comment_top, Inches(12.3), Inches(3.0))
        t = comment_box.text_frame
        t.text = "Комментарий ИИ — Динамика"
        t.paragraphs[0].font.size = Pt(14)
        t.paragraphs[0].font.name = self.settings.pptx_font_family
        t.paragraphs[0].font.bold = True
        try:
            # tighten inner margins for more space
            t.margin_left = Pt(2)
            t.margin_right = Pt(2)
            t.margin_top = Pt(2)
            t.margin_bottom = Pt(2)
        except Exception:
            pass
        prev_totals = prev
        cur_totals = cur
        ai_text = await self.gpt_service.generate_comparison_comment(prev_totals, cur_totals, "Динамика: предыдущий vs текущий")
        body = t.add_paragraph()
        body.text = ai_text
        body.font.size = Pt(11 if len(ai_text) > 600 else 12)
        body.font.name = self.settings.pptx_font_family
        try:
            for par in t.paragraphs:
                par.space_after = Pt(3)
            t.word_wrap = True
        except Exception:
            pass

    async def _add_top3_slide(self, prs: Presentation, period_data: Dict[str, ManagerData]) -> None:
        """Add slide with TOP-3 best and worst managers. Prefer AI ranking; fallback to metric-based."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        title = slide.shapes.title
        title.text = "ТОП‑3 лучших и ТОП‑3 худших"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.name = self.settings.pptx_font_family
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(204, 0, 0)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        try:
            title.left = 0
            title.width = prs.slide_width
        except Exception:
            pass

        # Try AI-based ranking
        ai_best: list[str] = []
        ai_worst: list[str] = []
        ai_reasons: dict[str, str] = {}

        try:
            kpi = {}
            for m in period_data.values():
                kpi[m.name] = {
                    'calls_plan': m.calls_plan,
                    'calls_fact': m.calls_fact,
                    'leads_units_plan': m.leads_units_plan,
                    'leads_units_fact': m.leads_units_fact,
                    'leads_volume_plan': m.leads_volume_plan,
                    'leads_volume_fact': m.leads_volume_fact,
                    'approved_volume': m.approved_volume,
                    'issued_volume': m.issued_volume,
                }
            ai = await self.gpt_service.rank_top3(kpi)
            ai_best = [n for n in ai.get('best', []) if n in period_data]
            ai_worst = [n for n in ai.get('worst', []) if n in period_data]
            ai_reasons = ai.get('reasons', {}) or {}
        except Exception:
            pass

        def fallback_top3() -> tuple[list[str], list[str]]:
            scored = []
            for m in period_data.values():
                score = 0.5 * (m.calls_percentage) + 0.5 * (m.leads_volume_percentage)
                scored.append((score, m.name))
            scored.sort(reverse=True)
            best_names = [name for _, name in scored[:3]]
            worst_names = [name for _, name in list(reversed(scored[-3:]))]
            return best_names, worst_names

        best_names = ai_best
        worst_names = ai_worst
        if not best_names or not worst_names:
            best_names, worst_names = fallback_top3()

        left = Inches(0.5)
        top = Inches(1.6)
        box_best = slide.shapes.add_textbox(left, top, Inches(6.0), Inches(4.5))
        tfb = box_best.text_frame
        tfb.text = "Лучшие:"
        tfb.paragraphs[0].font.name = self.settings.pptx_font_family
        tfb.paragraphs[0].font.size = Pt(20)
        for name in best_names:
            m = period_data[name]
            reason = ai_reasons.get(name, f"звонки {m.calls_percentage:.0f}%, объем {m.leads_volume_percentage:.0f}%")
            p = tfb.add_paragraph()
            p.text = f"🏆 {name}: {reason}"
            p.font.name = self.settings.pptx_font_family
            p.font.size = Pt(14)

        box_worst = slide.shapes.add_textbox(left + Inches(6.5), top, Inches(6.0), Inches(4.5))
        tfw = box_worst.text_frame
        tfw.text = "Ниже темпа:"
        tfw.paragraphs[0].font.name = self.settings.pptx_font_family
        tfw.paragraphs[0].font.size = Pt(20)
        for name in worst_names:
            m = period_data[name]
            reason = ai_reasons.get(name, f"звонки {m.calls_percentage:.0f}%, объем {m.leads_volume_percentage:.0f}%")
            p = tfw.add_paragraph()
            p.text = f"⚠️ {name}: {reason}"
            p.font.name = self.settings.pptx_font_family
            p.font.size = Pt(14)
    
    def _calculate_totals(self, period_data: Dict[str, ManagerData]) -> Dict[str, float]:
        """Calculate team totals."""
        totals = {
            'calls_plan': 0,
            'calls_fact': 0,
            'leads_units_plan': 0,
            'leads_units_fact': 0,
            'leads_volume_plan': 0.0,
            'leads_volume_fact': 0.0,
            'approved_volume': 0.0,
            'issued_volume': 0.0,
            'new_calls': 0,
            'new_calls_plan': 0,
        }
        
        for manager_data in period_data.values():
            totals['calls_plan'] += manager_data.calls_plan
            totals['calls_fact'] += manager_data.calls_fact
            totals['leads_units_plan'] += manager_data.leads_units_plan
            totals['leads_units_fact'] += manager_data.leads_units_fact
            totals['leads_volume_plan'] += manager_data.leads_volume_plan
            totals['leads_volume_fact'] += manager_data.leads_volume_fact
            totals['approved_volume'] += manager_data.approved_volume
            totals['issued_volume'] += manager_data.issued_volume
            totals['new_calls'] += manager_data.new_calls
            totals['new_calls_plan'] += manager_data.new_calls_plan
        
        # Calculate percentages
        totals['calls_percentage'] = (totals['calls_fact'] / totals['calls_plan'] * 100) if totals['calls_plan'] > 0 else 0
        totals['leads_units_percentage'] = (totals['leads_units_fact'] / totals['leads_units_plan'] * 100) if totals['leads_units_plan'] > 0 else 0
        totals['leads_volume_percentage'] = (totals['leads_volume_fact'] / totals['leads_volume_plan'] * 100) if totals['leads_volume_plan'] > 0 else 0
        
        return totals
    
    def _calculate_average_manager(self, period_data: Dict[str, ManagerData]) -> Dict[str, float]:
        """Calculate average manager baseline for Pro Core comparison."""
        if not period_data:
            return {}
        n = len(period_data)
        avg = {
            'calls_plan': 0,
            'calls_fact': 0,
            'leads_units_plan': 0,
            'leads_units_fact': 0,
            'leads_volume_plan': 0.0,
            'leads_volume_fact': 0.0,
            'approved_volume': 0.0,
            'issued_volume': 0.0,
            'new_calls': 0,
        }
        for m in period_data.values():
            avg['calls_plan'] += m.calls_plan
            avg['calls_fact'] += m.calls_fact
            avg['leads_units_plan'] += m.leads_units_plan
            avg['leads_units_fact'] += m.leads_units_fact
            avg['leads_volume_plan'] += m.leads_volume_plan
            avg['leads_volume_fact'] += m.leads_volume_fact
            avg['approved_volume'] += m.approved_volume
            avg['issued_volume'] += m.issued_volume
            avg['new_calls'] += m.new_calls
        for k in avg:
            avg[k] = avg[k] / n
        avg['calls_percentage'] = (avg['calls_fact'] / avg['calls_plan'] * 100) if avg['calls_plan'] > 0 else 0
        avg['leads_units_percentage'] = (avg['leads_units_fact'] / avg['leads_units_plan'] * 100) if avg['leads_units_plan'] > 0 else 0
        avg['leads_volume_percentage'] = (avg['leads_volume_fact'] / avg['leads_volume_plan'] * 100) if avg['leads_volume_plan'] > 0 else 0
        return avg
