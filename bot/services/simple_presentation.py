"""
Simplified PPTX generator matching client's reference exactly.
Structure:
- Title slide
- Manager statistics table (6 metrics x 5 columns)
- AI commentary (compact numbered list)
"""
from __future__ import annotations
import io
from typing import Dict
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm
from math import sqrt

# Plotly for charts
try:
    import plotly.graph_objects as go
    import plotly.io as pio
except Exception:
    go = None
    pio = None

from bot.config import Settings
from bot.services.yandex_gpt import YandexGPTService
from bot.services.data_aggregator import DataAggregatorService
from bot.services.di import Container


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


PRIMARY = "#1565C0"
ACCENT2 = "#42A5F5"
TEXT_MAIN = "#212121"
TEXT_MUTED = "#757575"
CARD_BG = "#F5F5F5"


class SimplePresentationService:
    """Generate simple PPTX presentation matching client's reference."""
    
    def __init__(self, settings: Settings):
        self.settings = settings
        self.ai = YandexGPTService(settings)
    
    async def generate_presentation(
        self,
        period_data: Dict,
        period_name: str,
        start_date: date,
        end_date: date,
        prev_data: Dict,
        prev_start: date,
        prev_end: date,
    ) -> bytes:
        """Generate presentation with title + one slide per manager + team summary."""
        prs = Presentation()
        prs.slide_width = Inches(11.69)  # 16:9
        prs.slide_height = Inches(6.58)
        
        margin = Inches(0.6)
        
        # Slide 1: Title
        await self._add_title_slide(prs, period_name, start_date, end_date, margin)
        
        # Calculate team averages for comparison
        total_managers = len(period_data)
        avg = self._calculate_averages(period_data, total_managers)
        prev_avg = self._calculate_averages(prev_data, len(prev_data)) if prev_data else None

        # Compute previous quarter references (weekly):
        prev_q_team_weekly, prev_q_per_manager_weekly = await self._compute_prev_quarter_refs(end_date)
        # Calls overview slide (second slide)
        await self._add_calls_overview_slide(prs, period_data, prev_data, prev_q_team_weekly, margin, period_name, start_date, end_date)
        # Leads overview slide (third slide)
        await self._add_leads_overview_slide(prs, period_data, prev_data, prev_q_team_weekly, margin, period_name, start_date, end_date)

        # Slide 4: Calls trend (line chart only)
        await self._add_calls_trend_slide(prs, period_data, margin, period_name, start_date, end_date)
        
        # Slide 5: TOP-2 leaders with detailed AI commentary
        await self._add_top2_leaders_slide(prs, period_data, prev_q_team_weekly, margin, period_name, start_date, end_date)
        
        # One slide per manager
        for manager_name, manager_data in period_data.items():
            await self._add_manager_stats_slide(prs, manager_name, manager_data, avg, prev_avg, prev_q_per_manager_weekly, margin, start_date, end_date)
        # Team summary slide removed per request
        
        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
    
    def _calculate_averages(self, period_data, total_managers):
        """Calculate team averages."""
        totals = {
            'calls_plan': sum(m.calls_plan for m in period_data.values()),
            'calls_fact': sum(m.calls_fact for m in period_data.values()),
            'new_calls_plan': sum(m.new_calls_plan for m in period_data.values()),
            'new_calls_fact': sum(m.new_calls for m in period_data.values()),
            'leads_units_plan': sum(m.leads_units_plan for m in period_data.values()),
            'leads_units_fact': sum(m.leads_units_fact for m in period_data.values()),
            'leads_volume_plan': sum(m.leads_volume_plan for m in period_data.values()),
            'leads_volume_fact': sum(m.leads_volume_fact for m in period_data.values()),
            'approved_units': sum(getattr(m, 'approved_units', 0) for m in period_data.values()),
            'issued_volume': sum(m.issued_volume for m in period_data.values()),
        }
        return {k: v / total_managers if total_managers > 0 else 0 for k, v in totals.items()}

    async def _compute_prev_quarter_refs(self, end_date):
        """Compute previous quarter DAILY averages over working days (exclude zero days).

        Returns team_daily_avg and per_manager_daily_avg dicts.
        """
        try:
            container = Container.get()
            aggregator = DataAggregatorService(container.sheets)
            from datetime import date as _date, timedelta
            # Determine previous quarter range based on provided end_date
            q = (end_date.month - 1) // 3 + 1
            cur_q_start_month = (q - 1) * 3 + 1
            cur_q_start = _date(end_date.year, cur_q_start_month, 1)
            prev_end = cur_q_start - timedelta(days=1)
            if q == 1:
                prev_q_year = end_date.year - 1
                prev_q_start_month = 10
            else:
                prev_q_year = end_date.year
                prev_q_start_month = (q - 2) * 3 + 1
            prev_start = _date(prev_q_year, prev_q_start_month, 1)
            series = await aggregator.get_daily_series(prev_start, prev_end)
            # Build daily lists (working days only), exclude zeros
            from datetime import datetime as _dt
            daily_values = {
                'calls_fact': [],
                'new_calls': [],
                'leads_units_fact': [],
                'leads_volume_fact': [],
                'approved_volume': [],
                'issued_volume': [],
            }
            for item in series:
                d = _dt.strptime(item['date'], '%Y-%m-%d').date()
                if d.weekday() >= 5:
                    continue
                for key in list(daily_values.keys()):
                    val = float(item.get(key, 0) or 0)
                    if val > 0:
                        daily_values[key].append(val)
            # Fallback: последние 20 рабочих дней до текущего периода
            if not any(daily_values.values()):
                alt_end = end_date - timedelta(days=1)
                alt_start = alt_end - timedelta(days=27)
                series = await aggregator.get_daily_series(alt_start, alt_end)
                daily_values = {k: [] for k in daily_values}
                for item in series:
                    d = _dt.strptime(item['date'], '%Y-%m-%d').date()
                    if d.weekday() >= 5:
                        continue
                    for key in list(daily_values.keys()):
                        val = float(item.get(key, 0) or 0)
                        if val > 0:
                            daily_values[key].append(val)
            def avg_of_key(key: str) -> float:
                arr = daily_values.get(key, [])
                if not arr:
                    return 0.0
                return sum(arr) / len(arr)
            team_daily = {
                'calls_fact': avg_of_key('calls_fact'),
                'new_calls_fact': avg_of_key('new_calls'),
                'leads_units_fact': avg_of_key('leads_units_fact'),
                'leads_volume_fact': avg_of_key('leads_volume_fact'),
                'approved_units': avg_of_key('approved_volume'),
                'issued_volume': avg_of_key('issued_volume'),
                'calls_plan': 0,
                'new_calls_plan': 0,
                'leads_units_plan': 0,
                'leads_volume_plan': 0.0,
            }
            # Per-manager daily average: approximate via team_daily / active managers
            prev_data = await aggregator._aggregate_data_for_period(prev_start, prev_end)
            active_managers = len(prev_data) if prev_data else 0
            per_manager_weekly = None
            if active_managers > 0:
                per_manager_weekly = {
                    'calls_fact': team_daily['calls_fact'] / active_managers,
                    'new_calls_fact': team_daily['new_calls_fact'] / active_managers,
                    'leads_units_fact': team_daily['leads_units_fact'] / active_managers,
                    'leads_volume_fact': team_daily['leads_volume_fact'] / active_managers,
                    'approved_units': team_daily['approved_units'] / active_managers,
                    'issued_volume': team_daily['issued_volume'] / active_managers,
                    'calls_plan': 0,
                    'new_calls_plan': 0,
                    'leads_units_plan': 0,
                    'leads_volume_plan': 0.0,
                }
            return team_daily, per_manager_weekly
        except Exception:
            return None, None

    async def _add_calls_overview_slide(self, prs, period_data, prev_data, prev_q_team_weekly, margin, period_name, start_date, end_date):
        """Add 'Общие показатели звонков' slide as per reference."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_logo(slide, prs)
        # Title
        title = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.5))
        tf = title.text_frame
        tf.text = "Общие показатели звонков"
        p = tf.paragraphs[0]; p.font.name = "Roboto"; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = hex_to_rgb(PRIMARY); p.alignment = PP_ALIGN.CENTER

        # Aggregate current totals
        cur_calls_plan = sum(m.calls_plan for m in period_data.values())
        cur_calls_fact = sum(m.calls_fact for m in period_data.values())
        cur_new_plan = sum(m.new_calls_plan for m in period_data.values())
        cur_new_fact = sum(m.new_calls for m in period_data.values())

        # Previous totals (fallback 0)
        prev_calls_fact = sum((getattr(m, 'calls_fact', 0) for m in (prev_data or {}).values())) if prev_data else 0
        prev_new_fact = sum((getattr(m, 'new_calls', 0) for m in (prev_data or {}).values())) if prev_data else 0
        prev_calls_plan = sum((getattr(m, 'calls_plan', 0) for m in (prev_data or {}).values())) if prev_data else 0
        prev_new_plan = sum((getattr(m, 'new_calls_plan', 0) for m in (prev_data or {}).values())) if prev_data else 0

        # Table with 2 rows + header: Показатель, План, Факт, Конверсия, % к факту (ПП), Δ конверсии, п.п. (ПП), Средний факт предыдущего квартала
        rows, cols = 3, 7
        tbl = slide.shapes.add_table(rows, cols, margin, Inches(0.9), prs.slide_width - 2*margin, Inches(1.4)).table
        headers = ["Показатель", "План", "Факт", "Конверсия", "% к факту (ПП)", "Δ конверсии, п.п. (ПП)", "Средний факт предыдущего квартала"]
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c); cell.text = h
            cell.fill.solid();
            # Подсветка правого блока колонок, относящихся к ПП (начиная с колонки 4)
            cell.fill.fore_color.rgb = hex_to_rgb("#E3F2FD" if c < 4 else "#BBDEFB")
            for par in cell.text_frame.paragraphs:
                par.font.name = "Roboto"; par.font.size = Pt(11); par.font.bold = True; par.alignment = PP_ALIGN.CENTER; par.font.color.rgb = hex_to_rgb(TEXT_MAIN)

        # Data rows
        def pct(a, b):
            try:
                return round(a / b * 100) if b else 0
            except Exception:
                return 0

        # Считаем дельту по факту и дельту конверсии (в п.п.) относительно ПП
        calls_conv = pct(cur_calls_fact, cur_calls_plan)
        new_conv = pct(cur_new_fact, cur_new_plan)
        prev_calls_conv = pct(prev_calls_fact, prev_calls_plan) if prev_calls_plan else 0
        prev_new_conv = pct(prev_new_fact, prev_new_plan) if prev_new_plan else 0
        vs_calls = pct(cur_calls_fact - prev_calls_fact, prev_calls_fact) if prev_calls_fact else 0
        vs_new = pct(cur_new_fact - prev_new_fact, prev_new_fact) if prev_new_fact else 0
        vs_calls_conv = calls_conv - prev_calls_conv
        vs_new_conv = new_conv - prev_new_conv

        # Средний факт предыдущего квартала, масштабированный под длительность текущего периода (рабочие дни)
        def workdays_between(a, b):
            cur = a
            days = 0
            from datetime import timedelta
            while cur <= b:
                if cur.weekday() < 5:
                    days += 1
                cur = cur + timedelta(days=1)
            return days
        wd = workdays_between(start_date, end_date)
        scale = (wd / 5) if wd else 0
        base_calls_week = (prev_q_team_weekly or {}).get('calls_fact', 0.0)
        base_new_week = (prev_q_team_weekly or {}).get('new_calls_fact', 0.0)
        sf_calls = base_calls_week * scale
        sf_new = base_new_week * scale
        data_rows = [
            ["Повторные звонки", cur_calls_plan, cur_calls_fact, f"{calls_conv}%", f"{vs_calls:+d}%", f"{vs_calls_conv:+d}", f"{sf_calls:.1f}"],
            ["Новые звонки", cur_new_plan, cur_new_fact, f"{new_conv}%", f"{vs_new:+d}%", f"{vs_new_conv:+d}", f"{sf_new:.1f}"],
        ]
        for r, row in enumerate(data_rows, start=1):
            for c, v in enumerate(row):
                cell = tbl.cell(r, c); cell.text = str(v)
                if r % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = hex_to_rgb("#F7F9FC")
                for par in cell.text_frame.paragraphs:
                    par.font.name = "Roboto"; par.font.size = Pt(10); par.alignment = PP_ALIGN.CENTER if c>0 else PP_ALIGN.LEFT
                    if c >= 4:  # колонки, относящиеся к ПП
                        par.font.color.rgb = hex_to_rgb(PRIMARY)

        # Примечание про ПП
        note = slide.shapes.add_textbox(margin, Inches(2.35), prs.slide_width - 2*margin, Inches(0.25))
        nt = note.text_frame; nt.clear()
        pnt = nt.paragraphs[0]; pnt.text = "Колонки справа (% к факту (ПП), Δ конверсии, п.п. (ПП), Средний факт предыдущего квартала) — сравнение с ПП."
        pnt.font.name = "Roboto"; pnt.font.size = Pt(8); pnt.font.color.rgb = hex_to_rgb(TEXT_MUTED)

        # Comment block
        box = slide.shapes.add_textbox(margin, Inches(2.65), prs.slide_width - 2*margin, Inches(3.6))
        prompt = (
            "Сформируй краткий комментарий по звонкам за период '" + period_name + "'. "
            f"Повторные: факт {cur_calls_fact} из {cur_calls_plan} ({calls_conv}%). Новые: факт {cur_new_fact} из {cur_new_plan} ({new_conv}%). "
            "Сравни с прошлым периодом, укажи изменения в %, сделай 2–3 тезиса и 1 рекомендацию."
        )
        try:
            text = await self.ai.generate_answer(prompt)
        except Exception:
            text = "Количество звонков выросло/снизилось относительно прошлого периода. Рекомендация: скорректировать темп и довести план."
        t = box.text_frame; t.clear();
        h = t.paragraphs[0]; h.text = "Комментарии нейросети"; h.font.name = "Roboto"; h.font.size = Pt(13); h.font.bold = True
        p1 = t.add_paragraph(); p1.text = text; p1.font.name = "Roboto"; p1.font.size = Pt(10)

    async def _add_leads_overview_slide(self, prs, period_data, prev_data, prev_q_team_weekly, margin, period_name, start_date, end_date):
        """Add 'Общие показатели по заявкам' slide (units, volume, approved, issued)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_logo(slide, prs)
        # Title
        title = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.5))
        tf = title.text_frame
        tf.text = "Общие показатели по заявкам"
        p = tf.paragraphs[0]; p.font.name = "Roboto"; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = hex_to_rgb(PRIMARY); p.alignment = PP_ALIGN.CENTER

        # Current totals
        units_plan = 0  # план по заявкам признан бессмысленным и исключён
        units_fact = sum(m.leads_units_fact for m in period_data.values())
        vol_plan = 0  # план по заявкам (млн) исключён
        vol_fact = sum(m.leads_volume_fact for m in period_data.values())
        approved_plan = sum(getattr(m, 'approved_plan', 0) for m in period_data.values())
        approved_fact = sum(getattr(m, 'approved_volume', 0) for m in period_data.values())
        issued_plan = sum(getattr(m, 'issued_plan', 0) for m in period_data.values())
        issued_fact = sum(m.issued_volume for m in period_data.values())

        # Previous totals
        prev_units_fact = sum((getattr(m, 'leads_units_fact', 0) for m in (prev_data or {}).values())) if prev_data else 0
        prev_vol_fact = sum((getattr(m, 'leads_volume_fact', 0.0) for m in (prev_data or {}).values())) if prev_data else 0
        prev_approved_fact = sum((getattr(m, 'approved_volume', 0.0) for m in (prev_data or {}).values())) if prev_data else 0
        prev_issued_fact = sum((getattr(m, 'issued_volume', 0.0) for m in (prev_data or {}).values())) if prev_data else 0

        # Table (compact): Показатель, Факт, % к факту (ПП), на средний факт предыдущего квартала
        rows, cols = 5, 4
        tbl = slide.shapes.add_table(rows, cols, margin, Inches(0.9), prs.slide_width - 2*margin, Inches(2.2)).table
        headers = ["Показатель", "Факт", "% к факту (ПП)", "Средний факт предыдущего квартала"]
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = hex_to_rgb("#E3F2FD" if c < 2 else "#BBDEFB")
            for par in cell.text_frame.paragraphs:
                par.font.name = "Roboto"; par.font.size = Pt(11); par.font.bold = True; par.alignment = PP_ALIGN.CENTER; par.font.color.rgb = hex_to_rgb(TEXT_MAIN)

        def pct(a, b):
            try:
                return round(a / b * 100) if b else 0
            except Exception:
                return 0

        # Колонки конверсии убраны

        # Масштабируем средние недельные значения прошлого квартала под рабочие дни текущего периода
        def workdays_between(a, b):
            cur = a
            days = 0
            from datetime import timedelta
            while cur <= b:
                if cur.weekday() < 5:
                    days += 1
                cur = cur + timedelta(days=1)
            return days
        wd = workdays_between(start_date, end_date)
        scale = (wd / 5) if wd else 0
        sf_units = ((prev_q_team_weekly or {}).get('leads_units_fact', 0.0)) * scale
        sf_vol = ((prev_q_team_weekly or {}).get('leads_volume_fact', 0.0)) * scale
        sf_appr = ((prev_q_team_weekly or {}).get('approved_units', 0.0)) * scale
        sf_iss = ((prev_q_team_weekly or {}).get('issued_volume', 0.0)) * scale
        data_rows = [
            ["Заявки, штук", units_fact, f"{pct(units_fact - prev_units_fact, prev_units_fact) if prev_units_fact else 0:+d}%", f"{sf_units:.1f}"],
            ["Заявки, млн", f"{vol_fact:.1f}", f"{pct(vol_fact - prev_vol_fact, prev_vol_fact) if prev_vol_fact else 0:+d}%", f"{sf_vol:.1f}"],
            ["Одобрено, млн", f"{approved_fact:.1f}", f"{pct(approved_fact - prev_approved_fact, prev_approved_fact) if prev_approved_fact else 0:+d}%", f"{sf_appr:.1f}"],
            ["Выдано, млн", f"{issued_fact:.1f}", f"{pct(issued_fact - prev_issued_fact, prev_issued_fact) if prev_issued_fact else 0:+d}%", f"{sf_iss:.1f}"],
        ]
        for r, row in enumerate(data_rows, start=1):
            for c, v in enumerate(row):
                cell = tbl.cell(r, c); cell.text = str(v)
                if r % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = hex_to_rgb("#F7F9FC")
                for par in cell.text_frame.paragraphs:
                    par.font.name = "Roboto"; par.font.size = Pt(10); par.alignment = PP_ALIGN.CENTER if c>0 else PP_ALIGN.LEFT
                    if c >= 2:
                        par.font.color.rgb = hex_to_rgb(PRIMARY)

        # Comment
        box = slide.shapes.add_textbox(margin, Inches(3.4), prs.slide_width - 2*margin, Inches(2.3))
        prompt = (
            "Сформируй комментарий по заявкам (шт и млн), одобрениям и выдачам за период '" + period_name + "'. "
            f"Факт: units {units_fact}/{units_plan}, volume {vol_fact:.1f}/{vol_plan:.1f}, approved {approved_fact:.1f}, issued {issued_fact:.1f}. "
            "Сравни с прошлым периодом, отметь сильные/слабые места, дай 2–3 рекомендации."
        )
        try:
            text = await self.ai.generate_answer(prompt)
        except Exception:
            text = "Факт заявок и объёмов оценён. Рекомендуется сфокусироваться на стабильности одобрений и доведении выдач."
        t = box.text_frame; t.clear();
        h = t.paragraphs[0]; h.text = "Комментарии нейросети"; h.font.name = "Roboto"; h.font.size = Pt(14); h.font.bold = True
        p1 = t.add_paragraph(); p1.text = text; p1.font.name = "Roboto"; p1.font.size = Pt(10)

    # Team summary slide removed per request
    
    async def _add_title_slide(self, prs, period_name, start_date, end_date, margin):
        """Title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_logo(slide, prs)
        
        # Title
        title_box = slide.shapes.add_textbox(
            margin, Inches(2.5), prs.slide_width - 2*margin, Inches(1.5)
        )
        title_box.text_frame.text = f"Отчет По Активности\n{period_name}"
        for p in title_box.text_frame.paragraphs:
            p.font.name = "Roboto"
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(PRIMARY)
            p.alignment = PP_ALIGN.CENTER
    
    async def _add_manager_stats_slide(self, prs, manager_name, manager_data, avg, prev_avg, prev_q_per_manager_weekly, margin, start_date, end_date):
        """Manager statistics table + AI commentary (exactly as in reference)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_logo(slide, prs)
        
        # Title with manager name
        title_box = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.5))
        title_box.text_frame.text = f"Статистика менеджера: {manager_name}"
        title_box.text_frame.paragraphs[0].font.name = "Roboto"
        title_box.text_frame.paragraphs[0].font.size = Pt(22)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Table: 8 rows (header + 7 metrics: добавлена "Одобрено, млн"), 6 columns
        rows = 8
        cols = 6
        tbl_width = prs.slide_width - 2*margin
        tbl_height = Inches(3.3)
        tbl = slide.shapes.add_table(rows, cols, margin, Inches(0.9), tbl_width, tbl_height).table
        
        col_widths = [Inches(2.2), Inches(1.6), Inches(1.5), Inches(1.7), Inches(1.8), Inches(2.0)]
        for c in range(6):
            tbl.columns[c].width = col_widths[c]
        
        headers = ["", "Запланировано", "Факт", "Конверсия", "Средний\nменеджер\nфакт", "Средний\nменеджер\nконверсия"]
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb("#E3F2FD")
            for p in cell.text_frame.paragraphs:
                p.font.name = "Roboto"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                p.alignment = PP_ALIGN.CENTER
        
        # Manager data vs averages
        m = manager_data
        # baseline for "средний менеджер": из прошлого квартала (недельный), если доступно; иначе текущая средняя
        # Базовый "средний менеджер" — недельный из прошлого квартала; масштабируем под рабочие дни текущего периода
        def workdays_between(a, b):
            cur = a
            days = 0
            from datetime import timedelta
            while cur <= b:
                if cur.weekday() < 5:
                    days += 1
                cur = cur + timedelta(days=1)
            return days
        wd = workdays_between(start_date, end_date)
        scale = (wd / 5) if wd else 0
        scaled_ref = None
        if prev_q_per_manager_weekly:
            scaled_ref = {k: (v * scale if isinstance(v, (int, float)) else v) for k, v in prev_q_per_manager_weekly.items()}

        ref = (scaled_ref or prev_avg or avg) or {
            'calls_fact':0,'new_calls_fact':0,'leads_units_fact':0,'leads_volume_fact':0.0,
            'approved_units':0,'issued_volume':0.0,'calls_plan':0,'new_calls_plan':0,'leads_units_plan':0,'leads_volume_plan':0.0
        }
        row_data = [
            ["Повторные звонки", m.calls_plan, m.calls_fact, 
             f"{(m.calls_fact/m.calls_plan*100) if m.calls_plan else 0:.0f}%",
             f"{ref['calls_fact']:.1f}",
             f"{(ref['calls_fact']/ref['calls_plan']*100) if ref['calls_plan'] else 0:.0f}%"],
            ["Новые звонки", m.new_calls_plan, m.new_calls,
             f"{(m.new_calls/m.new_calls_plan*100) if m.new_calls_plan else 0:.0f}%",
             f"{ref['new_calls_fact']:.1f}",
             f"{(ref['new_calls_fact']/ref['new_calls_plan']*100) if ref['new_calls_plan'] else 0:.0f}%"],
            ["Заведено заявок шт", "—", m.leads_units_fact,
             "—",
             f"{ref['leads_units_fact']:.1f}",
             f"{(ref['leads_units_fact']/ref['leads_units_plan']*100) if ref['leads_units_plan'] else 0:.0f}%"],
            ["Заведено заявок, млн", "—", f"{m.leads_volume_fact:.1f}",
             "—",
             f"{ref['leads_volume_fact']:.1f}",
             f"{(ref['leads_volume_fact']/ref['leads_volume_plan']*100) if ref['leads_volume_plan'] else 0:.0f}%"],
            ["Одобрено заявок шт", "", getattr(m, 'approved_units', 0), "", f"{ref['approved_units']:.1f}", ""],
            ["Одобрено, млн", "", f"{getattr(m, 'approved_volume', 0.0):.1f}", "", "", ""],
            ["Выдано, млн", "", f"{m.issued_volume:.1f}", "", f"{ref['issued_volume']:.1f}", ""]
        ]
        
        for r, data in enumerate(row_data, start=1):
            for c, val in enumerate(data):
                cell = tbl.cell(r, c)
                cell.text = str(val)
                try:
                    cell.vertical_anchor = 1  # Middle
                except Exception:
                    pass
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb("#F7F9FC")
                for p in cell.text_frame.paragraphs:
                    p.font.name = "Roboto"
                    p.font.size = Pt(10)
                    p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                    p.space_before = Pt(3)
                    p.space_after = Pt(3)
        
        # AI commentary section
        y_start = Inches(0.9) + tbl_height + Inches(0.2)
        comment_box = slide.shapes.add_textbox(margin, y_start, prs.slide_width - 2*margin, Inches(2.8))
        
        # Generate AI comment
        calls_conv = (m.calls_fact/m.calls_plan*100) if m.calls_plan else 0
        new_calls_conv = (m.new_calls/m.new_calls_plan*100) if m.new_calls_plan else 0
        leads_conv = (m.leads_units_fact/m.leads_units_plan*100) if m.leads_units_plan else 0
        vol_conv = (m.leads_volume_fact/m.leads_volume_plan*100) if m.leads_volume_plan else 0
        
        prompt = f"""Ты — руководитель отдела по банковским гарантиям. Проанализируй работу менеджера {manager_name} за период.

ПОКАЗАТЕЛИ МЕНЕДЖЕРА:
- Повторные звонки: {m.calls_fact} из {m.calls_plan} ({calls_conv:.0f}%) | Средний менеджер: {ref['calls_fact']:.1f}
- Новые звонки: {m.new_calls} из {m.new_calls_plan} ({new_calls_conv:.0f}%) | Средний менеджер: {ref['new_calls_fact']:.1f}
- Заявки (шт): {m.leads_units_fact} из {m.leads_units_plan} ({leads_conv:.0f}%) | Средний менеджер: {ref['leads_units_fact']:.1f}
- Заявки (млн): {m.leads_volume_fact:.1f} из {m.leads_volume_plan:.1f} ({vol_conv:.0f}%) | Средний менеджер: {ref['leads_volume_fact']:.1f}
- Одобрено: {getattr(m, 'approved_units', 0)} шт | Средний менеджер: {ref['approved_units']:.1f}
- Выдано (млн): {m.issued_volume:.1f} | Средний менеджер: {ref['issued_volume']:.1f}

ПРАВИЛА АНАЛИЗА:
1. Конверсия 80%+ = отлично, 60-80% = норма, <60% = нужна работа
2. Одобрено ≠ Выдано — это нормально (выдачи растягиваются по времени, не ругай конверсию)
3. Сравни КАЖДЫЙ показатель менеджера со средним по команде
4. Дай общий вывод: работает лучше/хуже/на уровне команды

ФОРМАТ ОТВЕТА (строго 3 пункта):
1. [Сравнение звонков и активности с командой — выше/ниже среднего, на сколько]
2. [Сравнение заявок и объёмов с командой — выше/ниже среднего, конверсия план/факт]
3. [Общий вывод и рекомендация — продолжить/усилить активность/проработать конверсию]

Пиши КРАТКО, ДЕЛОВЫМ ЯЗЫКОМ, БЕЗ ВОДЫ."""
        
        try:
            ai_text = await self.ai.generate_answer(prompt)
        except Exception as e:
            ai_text = (
                "1. Активность менеджера в работе\n"
                "2. Конверсия по заявкам и объёмам\n"
                "3. Общий вывод и рекомендации"
            )
        
        # Add title and text
        tf = comment_box.text_frame
        tf.clear()
        
        title_p = tf.paragraphs[0]
        title_p.text = "Комментарий ИИ"
        title_p.font.name = "Roboto"
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
        title_p.space_after = Pt(8)
        
        # Subtitle
        subtitle_p = tf.add_paragraph()
        subtitle_p.text = "Коротко и по делу:"
        subtitle_p.font.name = "Roboto"
        subtitle_p.font.size = Pt(11)
        subtitle_p.font.color.rgb = hex_to_rgb(TEXT_MUTED)
        subtitle_p.space_after = Pt(6)
        
        # AI text (compact)
        for line in ai_text.split('\n'):
            if line.strip():
                p = tf.add_paragraph()
                p.text = line.strip()
                p.font.name = "Roboto"
                p.font.size = Pt(10)
                p.font.color.rgb = hex_to_rgb(TEXT_MAIN)
                p.space_after = Pt(3)
                p.level = 0

    def _add_logo(self, slide, prs) -> None:
        """Add company logo at top-right if file exists (Логотип.png)."""
        import os
        from pptx.util import Inches
        logo_name_candidates = ["Логотип.png", "logo.png", "Logo.png"]
        root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
        logo_path = None
        for name in logo_name_candidates:
            p = os.path.join(root, name)
            if os.path.exists(p):
                logo_path = p
                break
        if not logo_path:
            return
        try:
            width = Inches(1.4)
            left = prs.slide_width - width - Inches(0.4)
            top = Inches(0.2)
            slide.shapes.add_picture(logo_path, left, top, width=width)
        except Exception:
            pass

    async def _add_calls_trend_slide(self, prs, period_data, margin, period_name, start_date, end_date):
        """Slide 4: Daily calls trend charts (team 2-line + manager multi-line)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_logo(slide, prs)
        # Title
        title = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.5))
        title.text_frame.text = "Динамика звонков"
        p = title.text_frame.paragraphs[0]; p.font.name = "Roboto"; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = hex_to_rgb(PRIMARY); p.alignment = PP_ALIGN.CENTER

        # Daily series
        container = Container.get()
        aggregator = DataAggregatorService(container.sheets)
        daily = await aggregator.get_daily_series(start_date, end_date)
        dates = [d['date'] for d in daily]
        calls = [d.get('calls_fact', 0) for d in daily]
        new_calls = [d.get('new_calls', 0) for d in daily]

        # Chart 1: Team (повторные + новые)
        chart1_top = Inches(0.9)
        chart1_h = Inches(2.5)
        try:
            import plotly.graph_objects as go
            import plotly.io as pio
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=dates, y=calls, mode='lines+markers', name='Повторные', line=dict(color='#1565C0', width=3)))
            fig.add_trace(go.Scatter(x=dates, y=new_calls, mode='lines+markers', name='Новые', line=dict(color='#42A5F5', width=3)))
            fig.update_layout(margin=dict(l=10,r=10,t=10,b=10), height=420, legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='right', x=1))
            from io import BytesIO
            stream = BytesIO(pio.to_image(fig, format='png', scale=2))
            slide.shapes.add_picture(stream, margin, chart1_top, width=prs.slide_width - 2*margin, height=chart1_h)
        except Exception:
            ph = slide.shapes.add_textbox(margin, chart1_top, prs.slide_width - 2*margin, chart1_h)
            ph.text_frame.text = "[График команды: установите plotly]"
            for par in ph.text_frame.paragraphs: par.font.size = Pt(12); par.alignment = PP_ALIGN.CENTER

        # Chart 2: Per-manager (sum new+repeat, each manager = line)
        chart2_top = Inches(3.5)
        chart2_h = Inches(2.7)
        try:
            from datetime import timedelta
            from collections import defaultdict
            daily_manager_buckets = defaultdict(lambda: defaultdict(lambda: {'calls':0,'new_calls':0}))
            all_records = container.sheets._reports.get_all_records()
            for rec in all_records:
                rec = {str(k).strip().lower(): v for k,v in rec.items()}
                date_str = str(rec.get('date','')).strip()
                if not date_str: continue
                d = aggregator._parse_record_date(date_str)
                if not d or d < start_date or d > end_date: continue
                mgr = str(rec.get('manager','')).strip()
                if not mgr: continue
                key = d.strftime('%Y-%m-%d')
                daily_manager_buckets[mgr][key]['calls'] += int(rec.get('evening_calls_success',0) or 0)
                daily_manager_buckets[mgr][key]['new_calls'] += int(rec.get('evening_new_calls',0) or 0)
            fig2 = go.Figure()
            colors_palette = ['#1565C0','#43A047','#E53935','#FB8C00','#8E24AA','#00ACC1','#FDD835','#7CB342']
            for i, (mgr, day_map) in enumerate(daily_manager_buckets.items()):
                vals = [day_map.get(d,{}).get('calls',0) + day_map.get(d,{}).get('new_calls',0) for d in dates]
                fig2.add_trace(go.Scatter(x=dates, y=vals, mode='lines+markers', name=mgr, line=dict(color=colors_palette[i % len(colors_palette)], width=2)))
            fig2.update_layout(margin=dict(l=10,r=10,t=10,b=10), height=460, legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='right', x=1))
            stream2 = BytesIO(pio.to_image(fig2, format='png', scale=2))
            slide.shapes.add_picture(stream2, margin, chart2_top, width=prs.slide_width - 2*margin, height=chart2_h)
        except Exception:
            ph = slide.shapes.add_textbox(margin, chart2_top, prs.slide_width - 2*margin, chart2_h)
            ph.text_frame.text = "[График по менеджерам: установите plotly или нет данных]"
            for par in ph.text_frame.paragraphs: par.font.size = Pt(10); par.alignment = PP_ALIGN.CENTER

    async def _add_top2_leaders_slide(self, prs, period_data, prev_q_team_weekly, margin, period_name, start_date, end_date):
        """Slide 5: TOP-2 rankings with detailed AI commentary."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_logo(slide, prs)
        title = slide.shapes.add_textbox(margin, Inches(0.3), prs.slide_width - 2*margin, Inches(0.5))
        title.text_frame.text = "Лидеры недели"
        p = title.text_frame.paragraphs[0]; p.font.name = "Roboto"; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = hex_to_rgb(PRIMARY); p.alignment = PP_ALIGN.CENTER

        # Calculate rankings
        from math import sqrt
        from pptx.enum.shapes import MSO_SHAPE
        # Baselines from prev quarter weekly scaled to current working days
        def workdays_between(a, b):
            cur=a; days=0
            from datetime import timedelta
            while cur<=b:
                if cur.weekday()<5: days+=1
                cur = cur + timedelta(days=1)
            return days
        wd = workdays_between(start_date, end_date)
        scale = (wd/5) if wd else 0
        base_new = (prev_q_team_weekly or {}).get('new_calls_fact', 0.0) * scale
        base_rep = (prev_q_team_weekly or {}).get('calls_fact', 0.0) * scale

        conv_raws = []
        temp = []
        for name, m in period_data.items():
            s_new = min(100, (m.new_calls/base_new*100) if base_new>0 else 0)
            s_rep = min(100, (m.calls_fact/base_rep*100) if base_rep>0 else 0)
            conv = (m.calls_fact/m.calls_plan*100) if m.calls_plan else 0
            raw = sqrt(max(0.0, conv) * max(0.0, float(m.calls_fact)))
            temp.append((name, s_new, s_rep, raw))
            conv_raws.append(raw)
        max_conv = max(conv_raws) if conv_raws else 1
        ranked_calls = []
        for name, s_new, s_rep, raw in temp:
            s_conv = (raw/max_conv*100) if max_conv>0 else 0
            total = s_new*0.4 + s_rep*0.3 + s_conv*0.3
            ranked_calls.append((name, round(total,1)))
        ranked_calls.sort(key=lambda x: x[1], reverse=True)
        top2_calls = ranked_calls[:2]

        # Issued weighted score 20/30/50
        issued_rank = []
        for name, m in period_data.items():
            score = (float(m.leads_volume_fact or 0)*0.2) + (float(getattr(m,'approved_volume',0) or 0)*0.3) + (float(m.issued_volume or 0)*0.5)
            issued_rank.append((name, round(score,1)))
        issued_rank.sort(key=lambda x: x[1], reverse=True)
        top2_issued = issued_rank[:2]

        # Guard: skip slide if no rankings
        if not top2_calls or not top2_issued:
            return
        
        # Render leaderboards as cards
        card_top = Inches(0.9)
        card_h = Inches(2.2)
        half_w = (prs.slide_width - margin*3) / 2

        # Left card: Calls
        left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin, card_top, half_w, card_h)
        left_card.fill.solid(); left_card.fill.fore_color.rgb = hex_to_rgb("#F0F8FF")
        left_card.line.color.rgb = hex_to_rgb(PRIMARY); left_card.line.width = Pt(1.5)
        ltf = left_card.text_frame; ltf.clear(); ltf.margin_left = Inches(0.2); ltf.margin_top = Inches(0.15)
        lh = ltf.paragraphs[0]; lh.text = "🏆 ТОП-2 по звонкам (40/30/30)"; lh.font.name = "Roboto"; lh.font.size = Pt(15); lh.font.bold = True; lh.font.color.rgb = hex_to_rgb(PRIMARY)
        for i, (name, score) in enumerate(top2_calls[:2], start=1):
            icon = "🥇" if i==1 else "🥈"
            lp = ltf.add_paragraph(); lp.text = f"{icon} {name}: {score}"; lp.font.name = "Roboto"; lp.font.size = Pt(14); lp.font.color.rgb = hex_to_rgb(TEXT_MAIN); lp.space_before = Pt(6)

        # Right card: Issued
        right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin + half_w + margin, card_top, half_w, card_h)
        right_card.fill.solid(); right_card.fill.fore_color.rgb = hex_to_rgb("#F1F8F4")
        right_card.line.color.rgb = hex_to_rgb("#43A047"); right_card.line.width = Pt(1.5)
        rtf = right_card.text_frame; rtf.clear(); rtf.margin_left = Inches(0.2); rtf.margin_top = Inches(0.15)
        rh = rtf.paragraphs[0]; rh.text = "💎 ТОП-2 по выданным (20/30/50)"; rh.font.name = "Roboto"; rh.font.size = Pt(15); rh.font.bold = True; rh.font.color.rgb = hex_to_rgb("#43A047")
        for i, (name, score) in enumerate(top2_issued[:2], start=1):
            icon = "🥇" if i==1 else "🥈"
            rp = rtf.add_paragraph(); rp.text = f"{icon} {name}: {score}"; rp.font.name = "Roboto"; rp.font.size = Pt(14); rp.font.color.rgb = hex_to_rgb(TEXT_MAIN); rp.space_before = Pt(6)

        # AI commentary below cards: explain formulas for both rankings
        comment_top = Inches(3.2)
        comment_h = Inches(3.0)
        cbox = slide.shapes.add_textbox(margin, comment_top, prs.slide_width - 2*margin, comment_h)
        ctf = cbox.text_frame; ctf.clear()
        ch = ctf.paragraphs[0]; ch.text = "Комментарии нейросети"; ch.font.name = "Roboto"; ch.font.size = Pt(14); ch.font.bold = True

        # Prepare data for AI prompt
        top_calls_data = []
        for name, score in top2_calls:
            m = period_data[name]
            s_new = min(100, (m.new_calls/base_new*100) if base_new>0 else 0)
            s_rep = min(100, (m.calls_fact/base_rep*100) if base_rep>0 else 0)
            conv = (m.calls_fact/m.calls_plan*100) if m.calls_plan else 0
            raw = sqrt(max(0.0, conv) * max(0.0, float(m.calls_fact)))
            s_conv = (raw/max_conv*100) if max_conv>0 else 0
            top_calls_data.append({
                'name': name, 'score': score, 
                'new': m.new_calls, 'repeat': m.calls_fact, 'plan': m.calls_plan, 'conv': conv,
                's_new': s_new, 's_rep': s_rep, 's_conv': s_conv
            })
        
        top_issued_data = []
        for name, score in top2_issued:
            m = period_data[name]
            top_issued_data.append({
                'name': name, 'score': score,
                'leads_vol': float(m.leads_volume_fact or 0),
                'approved': float(getattr(m,'approved_volume',0) or 0),
                'issued': float(m.issued_volume or 0)
            })

        # Build detailed prompt for AI with full formula breakdown
        # Guard against empty rankings
        if not top_calls_data or not top_issued_data:
            cp = ctf.add_paragraph(); cp.text = "Недостаточно данных для формирования рейтинга."; cp.font.name = "Roboto"; cp.font.size = Pt(10)
            return
        
        calls_1 = top_calls_data[0] if len(top_calls_data) > 0 else None
        calls_2 = top_calls_data[1] if len(top_calls_data) > 1 else None
        issued_1 = top_issued_data[0] if len(top_issued_data) > 0 else None
        issued_2 = top_issued_data[1] if len(top_issued_data) > 1 else None
        
        prompt_parts = [f"Сформируй краткие пояснения к рейтингам лидеров за {period_name}.", ""]
        prompt_parts.append("ТОП-2 ПО ЗВОНКАМ (веса: новые 40%, повторные 30%, конверсия 30%):")
        if calls_1:
            prompt_parts.append(f"🥇 {calls_1['name']}: {calls_1['score']} — новые {calls_1['new']}, повторные {calls_1['repeat']}, план {calls_1['plan']}, конв {calls_1['conv']:.0f}%")
            prompt_parts.append(f"   Балл новые {calls_1['s_new']:.1f}, балл повторные {calls_1['s_rep']:.1f}, балл конв {calls_1['s_conv']:.1f}")
            prompt_parts.append(f"   Расчёт: ({calls_1['s_new']:.1f}×0.4) + ({calls_1['s_rep']:.1f}×0.3) + ({calls_1['s_conv']:.1f}×0.3) = {calls_1['score']}")
        if calls_2:
            prompt_parts.append(f"🥈 {calls_2['name']}: {calls_2['score']} — новые {calls_2['new']}, повторные {calls_2['repeat']}, план {calls_2['plan']}, конв {calls_2['conv']:.0f}%")
            prompt_parts.append(f"   Балл новые {calls_2['s_new']:.1f}, балл повторные {calls_2['s_rep']:.1f}, балл конв {calls_2['s_conv']:.1f}")
            prompt_parts.append(f"   Расчёт: ({calls_2['s_new']:.1f}×0.4) + ({calls_2['s_rep']:.1f}×0.3) + ({calls_2['s_conv']:.1f}×0.3) = {calls_2['score']}")
        prompt_parts.append("")
        prompt_parts.append("ТОП-2 ПО ВЫДАННЫМ (веса: заведено 20%, одобрено 30%, выдано 50%):")
        if issued_1:
            prompt_parts.append(f"🥇 {issued_1['name']}: {issued_1['score']} — заведено {issued_1['leads_vol']:.1f} млн, одобрено {issued_1['approved']:.1f} млн, выдано {issued_1['issued']:.1f} млн")
            calc_1 = (issued_1['leads_vol']*0.2 + issued_1['approved']*0.3 + issued_1['issued']*0.5)
            prompt_parts.append(f"   Формула: ({issued_1['leads_vol']:.1f}×0.2) + ({issued_1['approved']:.1f}×0.3) + ({issued_1['issued']:.1f}×0.5) = {calc_1:.1f}")
        if issued_2:
            prompt_parts.append(f"🥈 {issued_2['name']}: {issued_2['score']} — заведено {issued_2['leads_vol']:.1f} млн, одобрено {issued_2['approved']:.1f} млн, выдано {issued_2['issued']:.1f} млн")
            calc_2 = (issued_2['leads_vol']*0.2 + issued_2['approved']*0.3 + issued_2['issued']*0.5)
            prompt_parts.append(f"   Формула: ({issued_2['leads_vol']:.1f}×0.2) + ({issued_2['approved']:.1f}×0.3) + ({issued_2['issued']:.1f}×0.5) = {calc_2:.1f}")
        prompt_parts.append("")
        prompt_parts.append("Дай краткий вывод (2–3 предложения): почему эти менеджеры лидеры и что команде стоит взять на заметку.")
        prompt = "\n".join(prompt_parts)
        
        try:
            comment = await self.ai.generate_answer(prompt)
        except Exception:
            comment = "Лидеры показали высокие результаты по взвешенным метрикам. Рекомендуется команде изучить их подход."
        cp = ctf.add_paragraph(); cp.text = comment; cp.font.name = "Roboto"; cp.font.size = Pt(10)

